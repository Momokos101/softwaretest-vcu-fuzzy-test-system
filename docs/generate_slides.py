from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

C_BG      = RGBColor(0x0D, 0x1B, 0x2A)
C_ACCENT  = RGBColor(0x00, 0xB4, 0xD8)
C_ACCENT2 = RGBColor(0x90, 0xE0, 0xEF)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT   = RGBColor(0xCA, 0xD3, 0xE0)
C_YELLOW  = RGBColor(0xFF, 0xD1, 0x66)
C_GREEN   = RGBColor(0x57, 0xCC, 0x99)
C_RED     = RGBColor(0xFF, 0x6B, 0x6B)
C_TABLE_H = RGBColor(0x00, 0x77, 0xA8)
C_ROW1    = RGBColor(0x12, 0x2A, 0x40)
C_ROW2    = RGBColor(0x0D, 0x1B, 0x2A)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── 工具函数 ──────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK)

def bg(slide):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = C_BG

def rect(slide, l, t, w, h, fc, lc=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fc
    if lc: s.line.color.rgb = lc
    else:  s.line.fill.background()
    return s

def box(slide, text, l, t, w, h,
        size=13, bold=False, color=C_WHITE,
        align=PP_ALIGN.LEFT):
    tf = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf.word_wrap = True
    p = tf.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return tf

def mbox(slide, lines, l, t, w, h, size=12, color=C_LIGHT, bold=False):
    """多行文本框，lines 是列表"""
    tf = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf.word_wrap = True
    frm = tf.text_frame; frm.word_wrap = True
    for i, line in enumerate(lines):
        p = frm.paragraphs[0] if i == 0 else frm.add_paragraph()
        p.space_before = Pt(3)
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold

def hline(slide, t, l=0.35, w=12.63):
    r = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(0.04))
    r.fill.solid(); r.fill.fore_color.rgb = C_ACCENT; r.line.fill.background()

def header(slide, title, sub=None):
    rect(slide, 0, 0, 13.33, 1.12, C_TABLE_H)
    box(slide, title, 0.4, 0.1, 10.5, 0.58, size=27, bold=True)
    if sub:
        box(slide, sub, 0.4, 0.68, 10.5, 0.38, size=13, color=C_ACCENT2)
    rect(slide, 11.75, 0.22, 1.25, 0.42, C_ACCENT)
    box(slide, "2026.05.21", 11.77, 0.24, 1.21, 0.36,
        size=11, align=PP_ALIGN.CENTER)

def table_header(slide, col_xs, col_ws, col_ts, y, rh=0.36):
    for cx, cw, ct in zip(col_xs, col_ws, col_ts):
        rect(slide, cx, y, cw, rh, C_TABLE_H)
        box(slide, ct, cx+0.06, y+0.07, cw-0.1, rh-0.1,
            size=12, bold=True, align=PP_ALIGN.CENTER)

def table_row(slide, col_xs, col_ws, texts, colors, y, rh=0.4, ri=0):
    bc = C_ROW1 if ri%2==0 else C_ROW2
    for cx, cw in zip(col_xs, col_ws):
        rect(slide, cx, y, cw, rh, bc)
    for cx, cw, txt, tc in zip(col_xs, col_ws, texts, colors):
        box(slide, txt, cx+0.07, y+0.08, cw-0.12, rh-0.1, size=11, color=tc)

# ══════════════════════════════════════════════════════════
# SLIDE 1  项目架构
# ══════════════════════════════════════════════════════════
s1 = add_slide(); bg(s1)
header(s1, "项目架构：两个系统，角色明确",
       "AutoTestDesign Tool（我们的工具）  ×  VCU 行为仿真器（被测目标应用）")

# 说明条
rect(s1, 0.35, 1.2, 12.63, 0.68, C_ROW1)
box(s1, "Assignment 要求：开发 AutoTestDesign Tool，选择一个目标应用，用工具去测它。",
    0.5, 1.24, 12.3, 0.28, size=13, color=C_ACCENT2)
box(s1, "我们的工具 → AutoTestDesign Tool（AI 驱动的测试设计系统）          "
        "我们的目标应用 → VCU 行为仿真器（Vehicle Control Unit 软件仿真）",
    0.5, 1.52, 12.3, 0.3, size=13, bold=True)

# 左卡：Tool
rect(s1, 0.35, 1.98, 0.48, 3.6, C_ACCENT)
box(s1, "T\nO\nO\nL", 0.37, 2.45, 0.44, 1.8,
    size=12, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
rect(s1, 0.83, 1.98, 5.65, 3.6, C_ROW1)
box(s1, "AutoTestDesign Tool", 0.95, 2.02, 5.4, 0.36,
    size=14, bold=True, color=C_ACCENT)
mbox(s1, [
    "● FR 1.0  从 CSV / 纯文本导入 VCU 需求",
    "● FR 1.1  解析 Input Fields、Data Ranges、Conditions、Expected Actions",
    "● FR 2.0  自动计算 Risk Score，标注 High / Medium / Low",
    "● FR 3.0  自动生成 EP / BVA / Decision Table 测试用例",
    "● FR 4.0  状态转移图建模，生成 All-States 覆盖序列  ★加分",
    "● FR 5.0  LLM 合成 Test Oracle  ★加分",
    "● FR 6.0  导出 JSON / CSV / Excel",
    "● FR 7.0  GAN 生成 CC2 时序序列，补充边界场景  ★加分",
    "● Interactive Review：每阶段支持人工修改",
], 0.95, 2.4, 5.45, 3.1, size=12)

# 箭头
box(s1, "HTTP POST\n/simulate\n→", 6.55, 3.25, 1.15, 0.9,
    size=12, color=C_YELLOW, align=PP_ALIGN.CENTER)

# 右卡：Simulator
rect(s1, 7.75, 1.98, 0.48, 3.6, C_ACCENT2)
box(s1, "S\nU\nT", 7.77, 2.55, 0.44, 1.4,
    size=12, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
rect(s1, 8.23, 1.98, 4.75, 3.6, C_ROW1)
box(s1, "VCU 行为仿真器", 8.35, 2.02, 4.5, 0.36,
    size=14, bold=True, color=C_ACCENT2)
mbox(s1, [
    "● 独立 FastAPI 服务（端口 8001）",
    "● 5 种输入信号，边界来自真实 BAIC HIL 数据（9615 条）",
    "● 状态机逻辑：IDLE → WAKING → READY → SLEEPING",
    "● 10 条明确需求（REQ-001 ~ REQ-010）",
    "● 每条需求有清晰的预期输出（Oracle）",
    "",
    "取代原北汽 API（已不可用）",
    "沿用汽车行业 SIL 测试惯例",
], 8.35, 2.4, 4.5, 3.1, size=12)

# 输出字段
rect(s1, 13.0, 2.5, 0.28, 2.2, C_ACCENT2)
box(s1, "O\nU\nT", 13.02, 2.72, 0.24, 1.2,
    size=10, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
mbox(s1, ["test_status", "vehicle_state", "vehicle_mode",
          "ready_flag", "detail"],
     13.02, 2.5, 0.26, 2.2, size=9, color=C_BG)

hline(s1, 5.7)
box(s1, "Risk Analysis Report、Test Plan、Detailed Test Design 均针对被测软件（VCU 仿真器），而非工具本身",
    0.35, 5.77, 12.63, 0.38, size=12, color=C_YELLOW)

# ══════════════════════════════════════════════════════════
# SLIDE 2  AutoTestDesign Tool FR 覆盖
# ══════════════════════════════════════════════════════════
s2 = add_slide(); bg(s2)
header(s2, "AutoTestDesign Tool：功能需求覆盖",
       "必须项 FR 1.0 / 1.1 / 2.0 / 3.0 / 6.0 全部实现；FR 4.0 / 5.0 / 7.0 作为加分项纳入")

fr_rows = [
    ("FR 1.0", "Input / Parsing",
     "从 CSV / 纯文本导入 VCU 需求，支持直接用户输入", "必须", C_GREEN),
    ("FR 1.1", "Requirement Structuring",
     "解析出 Input Fields、Data Ranges、Conditions、Expected Actions", "必须", C_GREEN),
    ("FR 2.0", "Risk Analysis",
     "对每条需求自动计算 Risk Score，输出 High / Medium / Low 优先级排序", "必须", C_GREEN),
    ("FR 3.0", "Black-Box Test Design",
     "自动生成 Equivalence Partitioning / Boundary Value Analysis / Decision Table 三种技术的用例", "必须", C_GREEN),
    ("FR 4.0", "White-Box Test Modeling",
     "状态转移图建模 VCU 状态机（IDLE/WAKING/READY/SLEEPING），生成 All-States 覆盖测试序列", "加分 ★", C_YELLOW),
    ("FR 5.0", "Test Oracle Generation",
     "LLM Prompt 根据需求描述 + 具体输入值，自动合成预期结果（Expected Result）", "加分 ★", C_YELLOW),
    ("FR 6.0", "Output & Export",
     "导出 JSON / CSV / Excel，含 Test Cases、Risk Scores、需求-用例追溯矩阵", "必须", C_GREEN),
    ("FR 7.0", "Test Suite Optimization",
     "GAN 模型生成 CC2 电压时序序列，补充 EP/BVA 无法覆盖的动态边界穿越场景", "加分 ★", C_YELLOW),
]

cxs = [0.35, 1.1, 3.38, 10.22]
cws = [0.72, 2.25, 6.8,  1.28]
table_header(s2, cxs, cws, ["ID", "功能类别", "我们的实现", "状态"], y=1.25)

for ri, (frid, cat, impl, st, sc) in enumerate(fr_rows):
    y = 1.61 + ri * 0.49
    table_row(s2, cxs, cws,
              [frid, cat, impl, st],
              [C_ACCENT, C_LIGHT, C_WHITE, sc],
              y=y, rh=0.47, ri=ri)

hline(s2, 6.22)
box(s2, "Interactive Review（Assignment「Mainly」核心要求）：",
    0.35, 6.3, 4.2, 0.32, size=13, bold=True, color=C_ACCENT)
steps = ["概念导入", "Coverage Item 识别", "策略选择",
         "用例 + 追溯", "Oracle 合成", "结果分析", "迭代改进"]
sx = 4.65
for i, st in enumerate(steps):
    rect(s2, sx, 6.28, 1.2, 0.34, C_TABLE_H)
    box(s2, st, sx+0.04, 6.3, 1.12, 0.28,
        size=11, align=PP_ALIGN.CENTER)
    sx += 1.22
    if i < len(steps)-1:
        box(s2, "→", sx-0.06, 6.3, 0.12, 0.28,
            size=13, color=C_ACCENT, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# SLIDE 3  VCU 仿真器
# ══════════════════════════════════════════════════════════
s3 = add_slide(); bg(s3)
header(s3, "VCU 行为仿真器：被测目标应用",
       "5 种输入信号  ×  状态机逻辑  ×  10 条需求  |  数据基础：真实 BAIC HIL 5 个数据库（9615 条）")

# 左：信号表
box(s3, "5 种输入信号及测试策略", 0.35, 1.22, 5.9, 0.32,
    size=13, bold=True, color=C_ACCENT)

scxs = [0.35, 1.88, 3.1, 4.35]
scws = [1.5,  1.2,  1.2, 1.85]
table_header(s3, scxs, scws, ["信号", "有效区间", "失效区间", "测试技术"], y=1.57)

sig_data = [
    ("CC2 电压",      "[4.8, 7.7]V",  "<4.8 或 >7.8V", "EP + BVA + 状态转移"),
    ("CC 电压值",     "区间外",        "[0.1, 3.9]V",   "EP + BVA"),
    ("CP 幅值",       "区间外",        "[9.1, 12.9]V",  "EP + BVA"),
    ("供电电压",      "区间外",        "[9.1, 15.9]V",  "EP + BVA"),
    ("网络唤醒使能",  "= 0",          "= 1",           "EP（二值穷举）"),
]
for ri, row in enumerate(sig_data):
    y = 1.93 + ri * 0.44
    table_row(s3, scxs, scws, row,
              [C_LIGHT, C_GREEN, C_RED, C_ACCENT2],
              y=y, rh=0.42, ri=ri)

# 10条需求
box(s3, "10 条需求（REQ-001 ~ REQ-010）", 0.35, 4.15, 5.9, 0.32,
    size=13, bold=True, color=C_ACCENT)
reqs = [
    ("REQ-001", "CC2 有效区间唤醒", C_GREEN),
    ("REQ-002", "CC2 低压 FAIL",   C_RED),
    ("REQ-003", "CC2 高压 FAIL",   C_RED),
    ("REQ-004", "休眠触发",         C_YELLOW),
    ("REQ-005", "CC 电压保护",      C_RED),
    ("REQ-006", "CP 幅值保护",      C_RED),
    ("REQ-007", "供电电压保护",     C_RED),
    ("REQ-008", "网络使能冲突",     C_RED),
    ("REQ-009", "字段一致性",       C_ACCENT2),
    ("REQ-010", "响应时序 ≤120s",  C_ACCENT2),
]
for i, (rid, rtitle, rc) in enumerate(reqs):
    col = i % 2; row = i // 2
    x = 0.35 + col * 3.1
    y = 4.5 + row * 0.45
    rect(s3, x, y, 2.95, 0.4, C_ROW1)
    rect(s3, x, y, 0.82, 0.4, rc)
    box(s3, rid, x+0.04, y+0.09, 0.74, 0.24,
        size=10, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    box(s3, rtitle, x+0.88, y+0.09, 2.0, 0.24, size=11, color=C_LIGHT)

# 右：状态机（三状态，均来自真实数据）
box(s3, "状态机逻辑（支撑 FR 4.0 白盒建模）", 6.55, 1.22, 6.4, 0.32,
    size=13, bold=True, color=C_ACCENT)

# 三个状态节点：READY 顶部居中，SLEEPING 左下，ERROR 右下
# READY (顶部中间)
rect(s3, 9.1, 1.62, 1.9, 0.52, C_GREEN)
box(s3, "READY", 9.1, 1.66, 1.9, 0.28,
    size=14, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
box(s3, "state=170  status=1", 9.1, 1.9, 1.9, 0.2,
    size=9, color=C_BG, align=PP_ALIGN.CENTER)

# SLEEPING (左下)
rect(s3, 6.9, 3.5, 1.9, 0.52, C_YELLOW)
box(s3, "SLEEPING", 6.9, 3.54, 1.9, 0.28,
    size=14, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
box(s3, "state=30  status=3", 6.9, 3.78, 1.9, 0.2,
    size=9, color=C_BG, align=PP_ALIGN.CENTER)

# ERROR (右下)
rect(s3, 11.2, 3.5, 1.9, 0.52, C_RED)
box(s3, "ERROR", 11.2, 3.54, 1.9, 0.28,
    size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
box(s3, "state=30  status=4", 11.2, 3.78, 1.9, 0.2,
    size=9, color=C_WHITE, align=PP_ALIGN.CENTER)

# 箭头标注（用文字替代真实箭头线）
# READY <-> SLEEPING
box(s3, "CC2=12.0V ->", 7.6, 2.45, 1.7, 0.24, size=10, color=C_YELLOW)
box(s3, "<- CC2 valid", 7.6, 2.72, 1.7, 0.24, size=10, color=C_GREEN)
# READY <-> ERROR
box(s3, "<- CC2 valid", 10.1, 2.45, 1.7, 0.24, size=10, color=C_GREEN)
box(s3, "无效信号 ->",  10.1, 2.72, 1.7, 0.24, size=10, color=C_RED)

# 状态转移表
box(s3, "状态转移规则（每条均有数据库记录对应）", 6.55, 4.2, 6.4, 0.3,
    size=12, bold=True, color=C_ACCENT)

tcxs = [6.55, 7.82, 9.75, 11.12]
tcws = [1.24, 1.9,  1.34, 1.8]
table_header(s3, tcxs, tcws,
             ["当前状态", "触发条件", "下一状态", "数据依据"],
             y=4.54, rh=0.33)

trans_rows = [
    ("任意",     "CC2 in [4.8, 7.7]V",  "READY",    "strategy=0, status=1, state=170"),
    ("READY",    "CC2 = 12.0V",         "SLEEPING", "strategy=-3, status=3, state=30"),
    ("任意",     "信号值在失效区间",     "ERROR",    "strategy=1, status=4, state=30"),
    ("SLEEPING", "CC2 in [4.8, 7.7]V",  "READY",    "strategy=0 重新唤醒"),
    ("ERROR",    "CC2 in [4.8, 7.7]V",  "READY",    "reset 后重新唤醒"),
]
sc_map = {"READY": C_GREEN, "SLEEPING": C_YELLOW, "ERROR": C_RED, "任意": C_LIGHT}
for ri, (cur, cond, nxt, note) in enumerate(trans_rows):
    y = 4.87 + ri * 0.37
    bc = C_ROW1 if ri % 2 == 0 else C_ROW2
    for cx, cw in zip(tcxs, tcws):
        rect(s3, cx, y, cw, 0.35, bc)
    box(s3, cur,  tcxs[0]+0.06, y+0.09, tcws[0]-0.1, 0.22,
        size=10, bold=True, color=sc_map.get(cur, C_LIGHT))
    box(s3, cond, tcxs[1]+0.06, y+0.09, tcws[1]-0.1, 0.22,
        size=10, color=C_LIGHT)
    box(s3, nxt,  tcxs[2]+0.06, y+0.09, tcws[2]-0.1, 0.22,
        size=10, bold=True, color=sc_map.get(nxt, C_LIGHT))
    box(s3, note, tcxs[3]+0.06, y+0.09, tcws[3]-0.1, 0.22,
        size=9, color=C_ACCENT2)

# ══════════════════════════════════════════════════════════
# SLIDE 4  测试计划 + Mainly 流程
# ══════════════════════════════════════════════════════════
s4 = add_slide(); bg(s4)
header(s4, "测试计划与测试设计流程（40% + 30% 交付物）",
       "Test Plan 覆盖 Assignment 1.2 节全部要求项  |  Detailed Test Design 遵循 Mainly 流程")

# ── 左半：测试套设计 ──
box(s4, "High-level Test Suite Design", 0.35, 1.22, 5.9, 0.3,
    size=13, bold=True, color=C_ACCENT)

scxs2 = [0.35, 1.18, 2.3,  3.45, 5.55]
scws2 = [0.8,  1.1,  1.12, 2.07, 0.75]
table_header(s4, scxs2, scws2,
             ["测试套", "需求", "信号/对象", "技术", "用例数"], y=1.55)

suites = [
    ("Suite A", "REQ-001~004", "CC2 电压",    "EP + BVA",             "~20"),
    ("Suite B", "REQ-005~008", "其余 4 信号", "EP + BVA",             "~16"),
    ("Suite C", "REQ-001,004", "状态机转移",  "State Transition（All States）", "~12"),
    ("Suite D", "REQ-001",     "GAN 时序",    "动态边界穿越序列",      "~10"),
    ("Suite E", "REQ-005~009", "多条件组合",  "Decision Table",        "~8"),
]
for ri, row in enumerate(suites):
    y = 1.91 + ri * 0.43
    table_row(s4, scxs2, scws2, row,
              [C_ACCENT, C_LIGHT, C_LIGHT, C_ACCENT2, C_YELLOW],
              y=y, rh=0.41, ri=ri)

rect(s4, 0.35, 4.08, 5.95, 0.34, C_TABLE_H)
box(s4, "合计：5 套，约 66 条用例，覆盖 ISO 29119-4 四种黑盒技术 + 白盒状态转移",
    0.45, 4.1, 5.8, 0.28, size=11, bold=True)

# Framework + Cost
box(s4, "Testing Framework & Cost", 0.35, 4.55, 4.0, 0.3,
    size=13, bold=True, color=C_ACCENT)
rect(s4, 0.35, 4.88, 5.95, 1.06, C_ROW1)
mbox(s4, [
    "框架：pytest + httpx  向 FastAPI 仿真器发 HTTP 请求执行测试用例",
    "选择理由：轻量、脚本化、可集成 CI、与 FastAPI 天然适配",
    "手工估算：66 用例 × 10 min ≈ 11 人时    工具自动执行：全套 < 5 分钟    节省 93%",
], 0.45, 4.92, 5.75, 0.96, size=12)

# ── 右半：Mainly 流程 ──
box(s4, "Detailed Test Design — Mainly 流程落地", 6.55, 1.22, 6.4, 0.3,
    size=13, bold=True, color=C_ACCENT)

flow = [
    ("① Concept 概念导入",
     "导入 REQ-001~010，解析 Input Fields / Data Ranges\nConditions / Expected Actions"),
    ("② Coverage Item 识别",
     "工具自动识别各信号等价类边界、状态机节点与转移边"),
    ("③ Coverage Strategy 选择",
     "按信号特征匹配技术：CC2→EP+BVA+状态转移 / 其余→EP+BVA+决策表\n支持人工修改策略"),
    ("④ Test Cases + Traceability",
     "生成用例，每条附 REQ-ID 追溯，输出需求-用例追溯矩阵"),
    ("⑤ Prompt Design / Oracle 合成",
     "LLM Prompt 模板根据需求描述+输入值合成预期输出（Expected Result）"),
    ("⑥ Results Analysis",
     "执行后统计 Pass/Fail，映射回 Coverage Item，展示覆盖率"),
    ("⑦ Improvement with Evidence",
     "对 Fail 用例：GAN 补充边界序列，新增 Coverage Item 形成闭环"),
]

for i, (title, body) in enumerate(flow):
    y = 1.56 + i * 0.73
    hc = C_TABLE_H if i < 4 else C_ACCENT
    rect(s4, 6.55, y, 2.0, 0.66, hc)
    box(s4, title, 6.62, y+0.1, 1.9, 0.46, size=11, bold=True, color=C_BG)
    rect(s4, 8.57, y, 4.38, 0.66, C_ROW1)
    mbox(s4, body.split("\n"), 8.65, y+0.08, 4.25, 0.52, size=11)
    if i < len(flow)-1:
        box(s4, "↓", 6.88, y+0.67, 0.3, 0.12,
            size=10, color=C_ACCENT, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# SLIDE 5  疑问
# ══════════════════════════════════════════════════════════
s5 = add_slide(); bg(s5)
header(s5, "我们的疑问——请老师指导",
       "以下问题希望在本次讨论中明确，以便调整后续方案")

questions = [
    ("Q1", "被测对象的选择是否合适？",
     "VCU 仿真器的核心逻辑是基于真实 HIL 数据归纳的条件判断，本质是规则驱动的 if-else 结构。\n"
     "Assignment 示例中列举了 web app、计算器、登录模块等目标。我们选择 VCU 仿真器是否符合要求？",
     C_ACCENT),
    ("Q2", "白盒测试的深度是否达标？",
     "VCU 仿真器逻辑较为直接，白盒测试（分支覆盖）用较少用例即可达到 100% 覆盖率。\n"
     "这样的白盒测试深度在评分上是否足够？还是需要被测软件具备更复杂的内部逻辑？",
     C_YELLOW),
    ("Q3", "非功能性测试如何处理？",
     "VCU 仿真器没有用户认证、并发场景、数据库交互等特征，非功能性测试目前仅覆盖响应时间（REQ-010）。\n"
     "对于这类嵌入式仿真目标，老师期望非功能性测试还应覆盖哪些方面？",
     C_RED),
    ("Q4", "GAN 模块如何归类？",
     "我们用 GAN 生成 CC2 电压时序序列，补充动态边界穿越场景的测试数据。\n"
     "这个模块对应 FR 7.0（Test Suite Optimization）还是 FR 3.0（Test Case Generation）？在评分中如何归类？",
     C_ACCENT2),
]

for i, (qid, qtitle, qbody, qc) in enumerate(questions):
    col = i % 2; row = i // 2
    x = 0.35 + col * 6.5
    y = 1.3  + row * 2.75
    rect(s5, x, y, 0.55, 2.6, qc)
    box(s5, qid, x+0.04, y+1.0, 0.47, 0.5,
        size=18, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    rect(s5, x+0.55, y, 5.6, 2.6, C_ROW1)
    box(s5, qtitle, x+0.65, y+0.14, 5.4, 0.38,
        size=14, bold=True, color=qc)
    mbox(s5, qbody.split("\n"), x+0.65, y+0.58, 5.38, 1.9, size=12)

# ── 保存 ─────────────────────────────────────────────────
path = "/Users/hby/Documents/GitHub/softwaretest-vcu-fuzzy-test-system/docs/teacher_discussion.pptx"
prs.save(path)
print("saved:", path)
