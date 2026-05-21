from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 色彩方案 ─────────────────────────────────────────────
C_BG      = RGBColor(0x0D, 0x1B, 0x2A)
C_ACCENT  = RGBColor(0x00, 0xB4, 0xD8)
C_ACCENT2 = RGBColor(0x90, 0xE0, 0xEF)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT   = RGBColor(0xCA, 0xD3, 0xE0)
C_YELLOW  = RGBColor(0xFF, 0xD1, 0x66)
C_GREEN   = RGBColor(0x57, 0xCC, 0x99)
C_RED     = RGBColor(0xFF, 0x6B, 0x6B)
C_TABLE_H = RGBColor(0x00, 0x77, 0xA8)
C_ROW1    = RGBColor(0x12, 0x2A, 0x40)
C_ROW2    = RGBColor(0x0D, 0x1B, 0x2A)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── 工具函数 ─────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK)

def bg(slide):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = C_BG

def rect(slide, l, t, w, h, fc, lc=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fc
    if lc: s.line.color.rgb = lc
    else:  s.line.fill.background()
    return s

def box(slide, text, l, t, w, h, size=13, bold=False,
        color=C_WHITE, align=PP_ALIGN.LEFT):
    tf = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf.word_wrap = True
    p = tf.text_frame.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return tf

def mbox(slide, lines, l, t, w, h, size=12, color=C_LIGHT, bold=False):
    tf = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf.word_wrap = True
    frm = tf.text_frame; frm.word_wrap = True
    for i, line in enumerate(lines):
        p = frm.paragraphs[0] if i == 0 else frm.add_paragraph()
        p.space_before = Pt(3)
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold

def hline(slide, t, l=0.35, w=12.63):
    r = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(0.04))
    r.fill.solid(); r.fill.fore_color.rgb = C_ACCENT; r.line.fill.background()

def header(slide, title, sub=None):
    rect(slide, 0, 0, 13.33, 1.12, C_TABLE_H)
    box(slide, title, 0.4, 0.1, 10.5, 0.58, size=27, bold=True)
    if sub:
        box(slide, sub, 0.4, 0.68, 10.5, 0.38, size=13, color=C_ACCENT2)
    rect(slide, 11.75, 0.22, 1.25, 0.42, C_ACCENT)
    box(slide, "2026.05.21", 11.77, 0.24, 1.21, 0.36,
        size=11, align=PP_ALIGN.CENTER)

def tbl_header(slide, cxs, cws, cts, y, rh=0.36):
    for cx, cw, ct in zip(cxs, cws, cts):
        rect(slide, cx, y, cw, rh, C_TABLE_H)
        box(slide, ct, cx+0.06, y+0.07, cw-0.1, rh-0.1,
            size=12, bold=True, align=PP_ALIGN.CENTER)

def tbl_row(slide, cxs, cws, texts, colors, y, rh=0.4, ri=0):
    bc = C_ROW1 if ri % 2 == 0 else C_ROW2
    for cx, cw in zip(cxs, cws):
        rect(slide, cx, y, cw, rh, bc)
    for cx, cw, txt, tc in zip(cxs, cws, texts, colors):
        box(slide, txt, cx+0.07, y+0.08, cw-0.12, rh-0.1, size=11, color=tc)

# ══════════════════════════════════════════════════════════
# SLIDE 1  项目架构
# ══════════════════════════════════════════════════════════
s1 = add_slide(); bg(s1)
header(s1, "项目架构：两个系统，角色明确",
       "AutoTestDesign Tool（我们的工具）  ×  VCU 行为仿真器（被测目标应用）")

# 说明条
rect(s1, 0.35, 1.2, 12.63, 0.62, C_ROW1)
box(s1, "Assignment 要求：开发 AutoTestDesign Tool，选择一个目标应用，用工具测它。",
    0.5, 1.24, 12.3, 0.26, size=13, color=C_ACCENT2)
box(s1, "我们的工具 → AutoTestDesign Tool          "
        "我们的目标应用 → VCU 行为仿真器（基于北汽真实 HIL 数据库重建，取代已不可用的原 API）",
    0.5, 1.5, 12.3, 0.26, size=12, bold=True)

# 左卡 Tool
rect(s1, 0.35, 1.92, 0.48, 3.75, C_ACCENT)
box(s1, "T\nO\nO\nL", 0.37, 2.5, 0.44, 1.8,
    size=11, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
rect(s1, 0.83, 1.92, 5.65, 3.75, C_ROW1)
box(s1, "AutoTestDesign Tool", 0.95, 1.96, 5.4, 0.34,
    size=14, bold=True, color=C_ACCENT)
mbox(s1, [
    "● FR 1.0/1.1  导入并解析 VCU 需求（CSV/文本）",
    "● FR 2.0       风险评估：High / Medium / Low",
    "● FR 3.0       黑盒用例：EP + BVA + 决策表 + 状态转移",
    "● FR 4.0 ★    白盒建模：state09/10/11 状态机，All-Transitions",
    "● FR 5.0 ★    LLM 合成 5 个输出字段的 Oracle",
    "● FR 6.0       导出 JSON（bq_new 格式）/ CSV / Excel",
    "● FR 7.0 ★    两轮模糊测试：DataVariation 三种变异策略",
    "● Interactive Review：每阶段支持人工修改",
], 0.95, 2.34, 5.45, 3.25, size=12)

# 箭头
box(s1, "HTTP\nPOST\n/simulate", 6.55, 3.2, 1.1, 0.9,
    size=11, color=C_YELLOW, align=PP_ALIGN.CENTER)

# 右卡 Simulator
rect(s1, 7.72, 1.92, 0.48, 3.75, C_ACCENT2)
box(s1, "S\nU\nT", 7.74, 2.55, 0.44, 1.4,
    size=11, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
rect(s1, 8.2, 1.92, 4.78, 3.75, C_ROW1)
box(s1, "VCU 行为仿真器", 8.32, 1.96, 4.5, 0.34,
    size=14, bold=True, color=C_ACCENT2)
mbox(s1, [
    "● 状态机：state09（休眠）→ state10（初始化）→ state11（运行）",
    "● 7路唤醒（w1-w7）+ 3路休眠条件（h1-h3）",
    "● 结果四分类：expected / error / stuck / new_state",
    "● 5个输出字段：vehicle_state / vehicle_mode /",
    "    power_current / bus_message_flag / pdcu_wake_reason",
    "● 已知缺陷：快速上下电序列触发 state10 卡死",
    "    → actual_duration 超时，来自北汽真实 bug 报告",
], 8.32, 2.34, 4.5, 3.25, size=12)

hline(s1, 5.8)
box(s1, "Risk Analysis Report、Test Plan、Detailed Test Design 均针对被测软件（VCU 仿真器），而非工具本身",
    0.35, 5.87, 12.63, 0.36, size=12, color=C_YELLOW)

# ══════════════════════════════════════════════════════════
# SLIDE 2  VCU 仿真器详细设计
# ══════════════════════════════════════════════════════════
s2 = add_slide(); bg(s2)
header(s2, "VCU 行为仿真器：被测目标应用",
       "状态机 + 卡死缺陷仿真 + 灰色地带设计（触发 new_state）")

# ── 左半 ────────────────────────────────────────────────

# 状态机
box(s2, "状态机", 0.35, 1.22, 2.0, 0.28, size=12, bold=True, color=C_ACCENT)
states = [("state09\n(休眠)", 0.4, 1.56, C_YELLOW),
          ("state10\n(初始化)", 2.3, 1.56, C_ACCENT),
          ("state11\n(运行)", 4.3, 1.56, C_GREEN)]
for name, lx, ly, fc in states:
    rect(s2, lx, ly, 1.7, 0.52, fc)
    box(s2, name, lx+0.05, ly+0.06, 1.6, 0.4,
        size=11, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
box(s2, "w任一->", 2.15, 1.68, 0.8, 0.22, size=9, color=C_LIGHT)
box(s2, "完成->",  4.05, 1.68, 0.72, 0.22, size=9, color=C_LIGHT)
box(s2, "<- h1+h2+h3", 2.65, 2.15, 1.5, 0.22, size=9, color=C_LIGHT)

# 输入信号表（压缩版）
box(s2, "输入信号", 0.35, 2.48, 2.0, 0.28,
    size=12, bold=True, color=C_ACCENT)
scxs = [0.35, 1.72, 2.98, 4.1]
scws = [1.34, 1.23, 1.09, 1.88]
tbl_header(s2, scxs, scws, ["信号","有效条件","时序","测试技术"], y=2.79, rh=0.3)
sig_rows = [
    ("w1 供电电压",   "> 9V",     "≥10ms", "EP+BVA(双轴)"),
    ("w2 CAN 报文",   "0x400~47F","无",    "EP"),
    ("w3 CP / w4 CC","> 9V/<4.4V","无",   "EP+BVA"),
    ("w5 CC2",        "UBR下降沿","无",    "EP+BVA"),
    ("w6/w7 口盖/门板",">4V/<1V", "≥10ms","EP+BVA(双轴)"),
    ("h1+h2+h3",      "三者同时=1","无",  "决策表(8组合)"),
]
for ri, row in enumerate(sig_rows):
    tbl_row(s2, scxs, scws, row,
            [C_LIGHT, C_GREEN, C_ACCENT2, C_ACCENT2],
            y=3.12+ri*0.36, rh=0.34, ri=ri)

# 卡死缺陷仿真逻辑
rect(s2, 0.35, 5.3, 5.88, 0.28, C_RED)
box(s2, "卡死缺陷仿真逻辑（Suite E 专项）", 0.42, 5.32, 5.75, 0.22,
    size=12, bold=True)
rect(s2, 0.35, 5.6, 5.88, 1.12, C_ROW1)
mbox(s2, [
    "仿真器维护最近唤醒时间戳列表，每次唤醒前检查：",
    "  若5秒内唤醒 ≥3次 且 所有相邻间隔 <1秒",
    "  → state=10（卡死），power_current>0.2A，duration=est_time×3",
    "  POST /reset 同时清空历史 → Suite A/B/C/F 不会意外触发",
], 0.45, 5.64, 5.72, 1.04, size=11)

# ── 右半 ────────────────────────────────────────────────

# 结果四分类
box(s2, "结果四分类", 6.55, 1.22, 3.0, 0.28,
    size=12, bold=True, color=C_ACCENT)
result_types = [
    ("expected",  "strategy=0",  "5字段全匹配预期", C_GREEN),
    ("error",     "strategy=1/2","状态未跳转，触发变异", C_YELLOW),
    ("stuck",     "strategy=3",  "state10卡死，duration×3", C_RED),
    ("new_state", "strategy=1/2","不匹配任何预期→灰色地带", C_ACCENT),
]
for i, (rt, st, desc, rc) in enumerate(result_types):
    y = 1.56 + i * 0.62
    rect(s2, 6.55, y, 1.45, 0.54, rc)
    box(s2, rt, 6.6, y+0.05, 1.36, 0.22, size=11, bold=True,
        color=C_BG, align=PP_ALIGN.CENTER)
    box(s2, st, 6.6, y+0.3,  1.36, 0.18, size=9,
        color=C_BG, align=PP_ALIGN.CENTER)
    rect(s2, 8.02, y, 5.0, 0.54, C_ROW1)
    box(s2, desc, 8.1, y+0.14, 4.85, 0.28, size=11, color=C_LIGHT)

# 灰色地带设计
rect(s2, 6.55, 4.08, 6.43, 0.28, C_ACCENT)
box(s2, "灰色地带设计（触发 new_state，使两轮模糊测试真正运转）",
    6.62, 4.1, 6.3, 0.22, size=11, bold=True, color=C_BG)
rect(s2, 6.55, 4.38, 6.43, 1.38, C_ROW1)
mbox(s2, [
    "问题：仿真器完全确定性，所有输出都被预定义，new_state 永远不会发生",
    "解决：在边界值处设置[灰色地带]——返回未预定义的输出组合",
    "示例：供电电压 = 9.0V（刚好等于边界）",
    "  → 返回 vehicle_state=10, power_current=0.08A（不属于正常也不属于卡死）",
    "  → Tool 匹配不上任何预期 → new_state → 第二轮围绕 9.0V 生成变异用例",
    "覆盖信号：w1(9.0V)、w3(9.0V)、w4(4.4V)、w6(4.0V)、w7(1.0V)",
], 6.63, 4.42, 6.3, 1.3, size=11)

# 执行机制（简化）
rect(s2, 6.55, 5.88, 6.43, 0.28, C_TABLE_H)
box(s2, "执行机制：单次请求，仿真器直接同步返回",
    6.62, 5.9, 6.3, 0.22, size=11, bold=True)
rect(s2, 6.55, 6.18, 6.43, 0.58, C_ROW1)
mbox(s2, [
    "POST /reset → POST /simulate → 返回5字段+result_type+actual_duration",
    "Tool 依次比较 expected → error → stuck，都不匹配则为 new_state",
], 6.63, 6.22, 6.3, 0.5, size=11)

# ══════════════════════════════════════════════════════════
# SLIDE 3  AutoTestDesign Tool
# ══════════════════════════════════════════════════════════
s3 = add_slide(); bg(s3)
header(s3, "AutoTestDesign Tool：功能需求覆盖",
       "必须项 FR1.0/1.1/2.0/3.0/6.0 全部实现；FR4.0/5.0/7.0 加分项已纳入")

fr_rows = [
    ("FR 1.0","Input/Parsing",
     "从 CSV/文本导入 REQ-001~014，字段：ID/Title/Input/Condition/Expected","必须",C_GREEN),
    ("FR 1.1","Requirement Structuring",
     "解析信号名/类型/范围、阈值/时序条件、5字段预期动作","必须",C_GREEN),
    ("FR 2.0","Risk Analysis",
     "按安全影响×缺陷概率打分，REQ-001/002/011/012 为 High（含真实 bug）","必须",C_GREEN),
    ("FR 3.0","Black-Box Test Design",
     "EP+BVA（电压+时序双轴）/ 决策表（h1×h2×h3 全8组合）/ 状态转移","必须",C_GREEN),
    ("FR 4.0","White-Box Test Modeling",
     "state09/10/11 状态机建模，All-States+All-Transitions，覆盖卡死路径","加分★",C_YELLOW),
    ("FR 5.0","Test Oracle Generation",
     "LLM Prompt 合成5个输出字段（vehicle_state/mode/current/flag/reason）的预期值","加分★",C_YELLOW),
    ("FR 6.0","Output & Export",
     "导出 JSON（与 bq_new assert 格式兼容）/ CSV / Excel，含需求-用例追溯矩阵","必须",C_GREEN),
    ("FR 7.0","Test Suite Optimization",
     "基于风险评分排序执行：High(REQ-001/002/011/012) → Medium → Low，时间不足时优先保障高风险用例","加分★",C_YELLOW),
]

cxs = [0.35, 1.1, 3.38, 10.22]
cws = [0.72, 2.25, 6.8, 1.28]
tbl_header(s3, cxs, cws, ["ID","功能类别","实现说明","状态"], y=1.25)
for ri, (fid, cat, impl, st, sc) in enumerate(fr_rows):
    y = 1.61 + ri * 0.49
    tbl_row(s3, cxs, cws, [fid,cat,impl,st],
            [C_ACCENT,C_LIGHT,C_WHITE,sc], y=y, rh=0.47, ri=ri)

# 补充：两轮模糊测试
rect(s3, 0.35, 5.58, 12.63, 0.52, C_ROW1)
box(s3, "补充（FR 3.0 扩展）：", 0.45, 5.62, 2.2, 0.28,
    size=12, bold=True, color=C_YELLOW)
box(s3,
    "两轮模糊测试（自适应测试，非 ISO 29119-4 标准技术）——第一轮运行82条用例，"
    "发现 new_state 后第二轮围绕灰色地带生成变异用例，属于测试用例生成的动态扩展手段",
    2.72, 5.62, 10.2, 0.44, size=12, color=C_LIGHT)

hline(s3, 6.22)
box(s3, "Interactive Review（Mainly 流程）：",
    0.35, 6.3, 3.5, 0.32, size=13, bold=True, color=C_ACCENT)
steps = ["需求导入","Coverage识别","策略选择","用例+追溯",
         "Oracle合成","执行+比较","结果分析","变异改进"]
sx = 4.05
for i, st in enumerate(steps):
    rect(s3, sx, 6.28, 1.15, 0.34, C_TABLE_H)
    box(s3, st, sx+0.04, 6.3, 1.07, 0.28,
        size=11, align=PP_ALIGN.CENTER)
    sx += 1.17
    if i < len(steps)-1:
        box(s3, "->", sx-0.06, 6.3, 0.1, 0.28,
            size=12, color=C_ACCENT, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# SLIDE 4  测试计划
# ══════════════════════════════════════════════════════════
s4 = add_slide(); bg(s4)
header(s4, "测试计划（40%）与测试设计流程（30%）",
       "6个测试套，82条用例  |  风险优先级排序执行  |  pytest + httpx 执行框架")

# 左：测试套表格
box(s4, "High-level Test Suite Design", 0.35, 1.22, 6.1, 0.3,
    size=13, bold=True, color=C_ACCENT)

scxs2 = [0.35, 1.2, 2.35, 3.62, 5.6]
scws2 = [0.82, 1.12, 1.24, 1.95, 0.72]
tbl_header(s4, scxs2, scws2, ["套件","需求","对象","技术","用例"], y=1.55)

suites = [
    ("Suite A","REQ-001~007","7路唤醒 EP",  "等价类（有效/无效/时序无效）","21"),
    ("Suite B","REQ-001/6/7","BVA 双轴",   "电压边界 × 时序边界","18"),
    ("Suite C","REQ-008~011","休眠决策表",  "h1×h2×h3 全8组合","8"),
    ("Suite D","REQ-001~012","状态转移",    "All-Transitions 5条路径","15"),
    ("Suite E","REQ-012",    "卡死缺陷",    "序列测试 + 重复执行变异","10"),
    ("Suite F","REQ-014",    "响应时序",    "actual_duration 趋势监控","10"),
]
for ri, row in enumerate(suites):
    tbl_row(s4, scxs2, scws2, row,
            [C_ACCENT,C_LIGHT,C_LIGHT,C_ACCENT2,C_YELLOW],
            y=1.91+ri*0.41, rh=0.39, ri=ri)

rect(s4, 0.35, 4.38, 6.1, 0.34, C_TABLE_H)
box(s4, "合计：6套，82条用例，覆盖 ISO 29119-4 四种黑盒技术 + 白盒状态转移",
    0.45, 4.4, 6.0, 0.28, size=11, bold=True)

# 框架与成本
box(s4, "Framework & Cost", 0.35, 4.85, 3.0, 0.3,
    size=13, bold=True, color=C_ACCENT)
rect(s4, 0.35, 5.18, 6.1, 1.58, C_ROW1)
mbox(s4, [
    "框架：pytest + httpx，向 FastAPI 仿真器发 HTTP 请求",
    "流程：POST /reset → POST /simulate → 仿真器同步返回5字段+duration+result_type",
    "Tool 比较返回值与预期 → 确定 strategy → 结果存 SQLite",
    "手工估算：82条 × 10min ≈ 14人时    工具自动执行：< 15分钟    节省 98%",
], 0.45, 5.22, 5.95, 1.48, size=11)

# 右：FR7.0 风险优先级 + 模糊测试补充
box(s4, "FR 7.0 实现：风险优先级排序执行", 6.55, 1.22, 6.4, 0.3,
    size=13, bold=True, color=C_ACCENT)

# 优先级排序图
priority_rows = [
    ("第一批", "High 风险", "REQ-001/002/011/012", C_RED,
     "唤醒核心路径 + 已知卡死缺陷，必须最先验证"),
    ("第二批", "Medium 风险", "REQ-003~005/013/014", C_YELLOW,
     "辅助唤醒信号 + 一致性 + 时序，主路径通过后执行"),
    ("第三批", "Low 风险",   "REQ-006/007/008/009", C_GREEN,
     "低优先级信号，时间充足时执行"),
]
for ri, (batch, level, reqs, bc, desc) in enumerate(priority_rows):
    y = 1.6 + ri * 0.88
    rect(s4, 6.55, y, 0.75, 0.78, bc)
    box(s4, batch, 6.58, y+0.24, 0.7, 0.28,
        size=11, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    rect(s4, 7.32, y, 1.35, 0.78, C_TABLE_H)
    box(s4, level, 7.35, y+0.03, 1.28, 0.28,
        size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    box(s4, reqs,  7.35, y+0.34, 1.28, 0.28,
        size=10, color=C_ACCENT2, align=PP_ALIGN.CENTER)
    rect(s4, 8.69, y, 4.26, 0.78, C_ROW1)
    box(s4, desc, 8.76, y+0.22, 4.14, 0.32, size=11, color=C_LIGHT)

box(s4, "时间不足时：第三批可跳过，优先保障 High 风险用例全部通过",
    6.55, 4.28, 6.4, 0.3, size=11, color=C_YELLOW)

# 模糊测试补充
rect(s4, 6.55, 4.65, 6.4, 0.28, C_ACCENT)
box(s4, "补充（FR 3.0 扩展）：两轮模糊测试——自适应测试，非 ISO 29119-4 标准技术",
    6.62, 4.67, 6.28, 0.22, size=11, bold=True, color=C_BG)
rect(s4, 6.55, 4.95, 6.4, 1.8, C_ROW1)
mbox(s4, [
    "第一轮（82条）：按高→中→低风险顺序执行，new_state 和 stuck 写入数据库",
    "第二轮（动态）：查询第一轮 new_state → based_new_state_fuzz 围绕灰色地带生成变异用例",
    "   示例：9.0V 触发 new_state → 第二轮自动生成 8.8V / 8.9V / 9.1V / 9.2V 四条变异用例",
    "作用：FR 3.0 的动态扩展——不是排序优化，而是基于覆盖发现补充遗漏边界",
], 6.63, 4.99, 6.28, 1.7, size=11)

# ══════════════════════════════════════════════════════════
# SLIDE 5  详细测试设计
# ══════════════════════════════════════════════════════════
s5 = add_slide(); bg(s5)
header(s5, "详细测试设计：黑盒 + 白盒技术",
       "针对 VCU 仿真器  |  ISO 29119-4 四种黑盒技术 + 白盒 All-Transitions 覆盖")

# ── 左半：黑盒 ──────────────────────────────────────────
box(s5, "黑盒测试技术", 0.35, 1.22, 6.1, 0.32,
    size=14, bold=True, color=C_ACCENT)

# 1. EP
rect(s5, 0.35, 1.58, 6.1, 0.3, C_TABLE_H)
box(s5, "① 等价类划分（EP）— Suite A：7路唤醒信号，每路3个等价类",
    0.45, 1.6, 6.0, 0.24, size=12, bold=True)

ep_cxs = [0.35, 1.52, 2.75, 4.12]
ep_cws = [1.14, 1.2,  1.34, 2.05]
tbl_header(s5, ep_cxs, ep_cws, ["信号","有效类","无效类(值)","无效类(时序)"], y=1.91, rh=0.3)
ep_rows = [
    ("w1 供电电压", "9.3V, 15ms", "8.5V（值不足）", "9.3V, 8ms（时序不足）"),
    ("w4 CC 电压",  "3.0V",       "5.0V（超出范围）","—（无时序要求）"),
    ("w2 CAN报文",  "0x420",      "0x300（超范围）", "—"),
]
for ri, row in enumerate(ep_rows):
    tbl_row(s5, ep_cxs, ep_cws, row,
            [C_LIGHT, C_GREEN, C_RED, C_YELLOW],
            y=2.21+ri*0.35, rh=0.33, ri=ri)

# 2. BVA
rect(s5, 0.35, 3.3, 6.1, 0.3, C_TABLE_H)
box(s5, "② 边界值分析（BVA）— Suite B：电压 × 时序 双轴边界",
    0.45, 3.32, 6.0, 0.24, size=12, bold=True)

bva_cxs = [0.35, 1.85, 3.25, 4.65]
bva_cws = [1.47, 1.37, 1.37, 1.41]
tbl_header(s5, bva_cxs, bva_cws, ["边界类型","下边界-1","下边界","下边界+1"], y=3.63, rh=0.3)
bva_rows = [
    ("w1 电压轴", "8.9V → error",  "9.0V → error", "9.1V → expected"),
    ("w1 时序轴", "9ms → error",   "10ms → expected","11ms → expected"),
    ("w6 电压轴", "3.9V → error",  "4.0V → error", "4.1V → expected"),
]
for ri, row in enumerate(bva_rows):
    tbl_row(s5, bva_cxs, bva_cws, row,
            [C_ACCENT2, C_RED, C_YELLOW, C_GREEN],
            y=3.96+ri*0.35, rh=0.33, ri=ri)

# 3. 决策表
rect(s5, 0.35, 5.08, 6.1, 0.3, C_TABLE_H)
box(s5, "③ 决策表（Decision Table）— Suite C：休眠条件 h1 AND h2 AND h3（8组合）",
    0.45, 5.1, 6.0, 0.24, size=12, bold=True)
dt_cxs = [0.35, 1.05, 1.75, 2.45, 3.15]
dt_cws = [0.67, 0.67, 0.67, 0.67, 2.9]
tbl_header(s5, dt_cxs, dt_cws, ["h1","h2","h3","×4组","预期结果"], y=5.41, rh=0.28)
for ri, (h1,h2,h3,cnt,exp,clr) in enumerate([
    ("0","×","×","(1~4组)","维持 state11，不休眠",C_RED),
    ("1","1","1","(第8组)","state09，bus_flag=0",C_GREEN),
]):
    y = 5.69 + ri * 0.34
    bc = C_ROW1 if ri%2==0 else C_ROW2
    for cx,cw in zip(dt_cxs,dt_cws): rect(s5,cx,y,cw,0.32,bc)
    for cx,cw,txt,tc in zip(dt_cxs,dt_cws,[h1,h2,h3,cnt,exp],[C_LIGHT,C_LIGHT,C_LIGHT,C_ACCENT2,clr]):
        box(s5,txt,cx+0.05,y+0.07,cw-0.08,0.22,size=10,color=tc)

# ── 右半：白盒 ─────────────────────────────────────────
box(s5, "白盒测试技术（FR 4.0）", 6.6, 1.22, 6.3, 0.32,
    size=14, bold=True, color=C_ACCENT)

# 状态转移图
rect(s5, 6.6, 1.58, 6.3, 0.3, C_TABLE_H)
box(s5, "④ 状态转移测试（State Transition）— Suite D + E：All-Transitions",
    6.7, 1.6, 6.2, 0.24, size=12, bold=True)

# 状态节点
for name, lx, ly, fc in [
    ("state09\n(休眠)", 6.85, 2.0,  C_YELLOW),
    ("state10\n(初始化)",8.85, 2.0, C_ACCENT),
    ("state11\n(运行)", 10.85, 2.0, C_GREEN),
]:
    rect(s5, lx, ly, 1.6, 0.52, fc)
    box(s5, name, lx+0.05, ly+0.06, 1.5, 0.4,
        size=12, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
box(s5, "w有效->", 8.5,  2.12, 0.82, 0.22, size=10, color=C_LIGHT)
box(s5, "完成->",  10.5, 2.12, 0.72, 0.22, size=10, color=C_LIGHT)

# 转移路径表
rect(s5, 6.6, 2.63, 6.3, 0.3, C_TABLE_H)
box(s5, "5条转移路径（All-Transitions 覆盖准则）",
    6.7, 2.65, 6.2, 0.24, size=12, bold=True)

tr_cxs = [6.6,  7.6,  9.0,  10.4]
tr_cws = [0.97, 1.37, 1.37, 2.22]
tbl_header(s5, tr_cxs, tr_cws,
           ["路径", "起始状态", "终止状态", "触发条件 / 预期结果"],
           y=2.96, rh=0.3)
tr_rows = [
    ("T-001", "state09", "state11", "w1有效(9.3V,15ms) → expected"),
    ("T-002", "state11", "state09", "h1+h2+h3全=1 → expected，bus_flag=0"),
    ("T-003", "state11", "state11", "仅h1+h2=1（缺h3）→ error，不休眠"),
    ("T-004", "state09", "state09", "w1无效(8.5V) → error，state不变"),
    ("T-005", "state09", "state10", "连续3次唤醒间隔<1s → stuck，duration超时"),
]
tr_colors = [C_GREEN, C_GREEN, C_YELLOW, C_RED, C_RED]
for ri, (pid,src,dst,cond) in enumerate(tr_rows):
    y = 3.29 + ri * 0.39
    tbl_row(s5, tr_cxs, tr_cws, [pid,src,dst,cond],
            [C_ACCENT, C_ACCENT2, tr_colors[ri], C_LIGHT],
            y=y, rh=0.37, ri=ri)

# 白盒覆盖目标（压缩，给非功能测试留空间）
rect(s5, 6.6, 5.28, 6.3, 0.28, C_TABLE_H)
box(s5, "白盒覆盖目标汇总",
    6.7, 5.3, 6.2, 0.22, size=11, bold=True)

wb_cxs = [6.6, 8.35, 10.9]
wb_cws = [1.72, 2.52, 1.72]
tbl_header(s5, wb_cxs, wb_cws, ["覆盖准则","对应路径","Suite"], y=5.58, rh=0.26)
for ri, (cr, wbpath, suite, clr) in enumerate([
    ("All-States",     "state09/10/11 三状态全覆盖", "Suite D",   C_GREEN),
    ("All-Transitions","T-001~T-005 五条路径全覆盖", "Suite D+E", C_GREEN),
    ("Branch Coverage","7路唤醒 × 有效/无效 分支",   "Suite A",   C_ACCENT2),
    ("Stuck Path",     "T-005 卡死缺陷专项路径",     "Suite E",   C_RED),
]):
    y = 5.84 + ri * 0.3
    bc = C_ROW1 if ri%2==0 else C_ROW2
    for cx,cw in zip(wb_cxs,wb_cws): rect(s5,cx,y,cw,0.28,bc)
    for cx,cw,txt,tc in zip(wb_cxs,wb_cws,[cr,wbpath,suite],[clr,C_LIGHT,C_ACCENT]):
        box(s5,txt,cx+0.05,y+0.06,cw-0.08,0.2,size=10,color=tc)

# 非功能测试（底部横跨全页）
hline(s5, 7.06, l=0.35, w=12.63)
rect(s5, 0.35, 7.08, 2.4, 0.3, C_TABLE_H)
box(s5, "非功能测试（Suite F）", 0.42, 7.1, 2.28, 0.24,
    size=11, bold=True)
rect(s5, 2.77, 7.08, 10.2, 0.3, C_ROW1)
mbox(s5, [
    "REQ-014：actual_duration ≤ 20s（type1）/ ≤ 60s（type2）；"
    "跨10次分析 duration 趋势，检测性能退化；"
    "power_current > 0.2A 为功耗异常指标（来自北汽原始 bug 描述）"
], 2.84, 7.1, 10.1, 0.26, size=11)

# ══════════════════════════════════════════════════════════
# SLIDE 6  疑问
# ══════════════════════════════════════════════════════════
s6 = add_slide(); bg(s6)
header(s6, "我们的疑问——请老师指导",
       "以下问题希望在本次讨论中明确，以便调整后续方案")

questions = [
    ("Q1", "被测对象的选择是否合适？",
     "VCU 仿真器基于北汽真实 HIL 数据库重建，包含真实记录的 state10 卡死缺陷（0319文档）。\n"
     "与 web app 等相比，这类嵌入式 SIL 仿真目标是否符合 Assignment 对「自选目标应用」的要求？",
     C_ACCENT),
    ("Q2", "人造 stuck bug 对白盒测试的价值老师如何评价？",
     "我们在仿真器中主动模拟了 state10 卡死缺陷（来源于真实 bug 报告），Suite E 专门针对此路径设计。\n"
     "老师认为这种「有意引入已知缺陷供测试发现」的方式，对白盒路径覆盖的评分是否有效？",
     C_YELLOW),
    ("Q3", "两轮模糊测试反馈机制对应 FR 7.0 吗？",
     "我们实现了：第一轮发现 new_state → 存库 → 第二轮 based_new_state_fuzz 变异再测。\n"
     "这是否满足 FR 7.0「Test Suite Optimization based on risk or coverage efficiency」的要求？",
     C_RED),
    ("Q4", "GAN 在项目中如何归类，对评分有贡献吗？",
     "GAN 生成 CC2 电压时序数据，它是测试数据生成工具，不是 ISO 29119-4 定义的黑盒/白盒技术。\n"
     "请问 GAN 应归入 FR 3.0（测试用例生成）还是 FR 7.0（测试套优化）？对哪个评分项有实质贡献？",
     C_ACCENT2),
]

for i, (qid, qtitle, qbody, qc) in enumerate(questions):
    col = i % 2; row = i // 2
    x = 0.35 + col * 6.5
    y = 1.3  + row * 2.75
    rect(s6, x, y, 0.55, 2.6, qc)
    box(s6, qid, x+0.04, y+1.0, 0.47, 0.5,
        size=18, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    rect(s6, x+0.55, y, 5.6, 2.6, C_ROW1)
    box(s6, qtitle, x+0.65, y+0.14, 5.4, 0.38,
        size=14, bold=True, color=qc)
    mbox(s6, qbody.split("\n"), x+0.65, y+0.58, 5.38, 1.9, size=12)

# ── 保存 ─────────────────────────────────────────────────
path = "/Users/hby/Documents/GitHub/softwaretest-vcu-fuzzy-test-system/docs/teacher_discussion_v7.pptx"
prs.save(path)
print("saved:", path)
