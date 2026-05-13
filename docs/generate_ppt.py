"""
AutoTestDesign 期末项目启动会 PPT 生成脚本
运行方式: python3 docs/generate_ppt.py
输出: docs/AutoTestDesign_启动会.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
import copy

# ── 颜色定义 ──────────────────────────────────────────────
BLUE_DARK   = RGBColor(0x1E, 0x40, 0xAF)   # #1E40AF 深蓝
BLUE_MID    = RGBColor(0x29, 0x63, 0xD0)   # 中蓝
BLUE_LIGHT  = RGBColor(0xDB, 0xEA, 0xFE)   # #DBEAFE 浅蓝背景
SLATE_900   = RGBColor(0x0F, 0x17, 0x2A)   # 近黑
SLATE_700   = RGBColor(0x33, 0x4D, 0x6B)   # 深灰蓝
SLATE_500   = RGBColor(0x64, 0x74, 0x8B)   # 中灰
SLATE_200   = RGBColor(0xE2, 0xE8, 0xF0)   # 浅灰
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GREEN       = RGBColor(0x16, 0xA3, 0x4A)
RED         = RGBColor(0xDC, 0x26, 0x26)
ORANGE      = RGBColor(0xEA, 0x58, 0x0C)
YELLOW_BG   = RGBColor(0xFF, 0xF7, 0xED)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── 工具函数 ─────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # 完全空白布局
    return prs.slides.add_slide(layout)


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width) if line_width else Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=SLATE_900,
             align=PP_ALIGN.LEFT, font_name="Microsoft YaHei",
             italic=False, word_wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_textbox_lines(slide, lines, left, top, width, height,
                      font_size=16, color=SLATE_900, bold=False,
                      line_spacing=1.2, font_name="Microsoft YaHei",
                      align=PP_ALIGN.LEFT):
    """多行文本框，lines 为 list of (text, bold, color, size)"""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            txt, b, c, sz = item, bold, color, font_size
        else:
            txt = item[0]
            b   = item[1] if len(item) > 1 else bold
            c   = item[2] if len(item) > 2 else color
            sz  = item[3] if len(item) > 3 else font_size
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(sz)
        run.font.bold = b
        run.font.color.rgb = c
        run.font.name = font_name
    return txBox


def header_bar(slide, title, subtitle=None):
    """顶部深蓝标题条"""
    add_rect(slide, 0, 0, 13.33, 1.15, fill_color=BLUE_DARK)
    add_text(slide, title, 0.4, 0.15, 10, 0.7,
             font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.75, 10, 0.4,
                 font_size=14, color=RGBColor(0xBF, 0xDB, 0xFF),
                 align=PP_ALIGN.LEFT)
    # 底部细线装饰
    add_rect(slide, 0, 1.15, 13.33, 0.04, fill_color=BLUE_MID)


def footer_bar(slide, text="AutoTestDesign 启动会 · 2026"):
    add_rect(slide, 0, 7.2, 13.33, 0.3, fill_color=SLATE_200)
    add_text(slide, text, 0.3, 7.22, 12.7, 0.25,
             font_size=10, color=SLATE_500, align=PP_ALIGN.LEFT)


def tag(slide, text, left, top, width=1.4, height=0.35,
        bg=BLUE_LIGHT, fg=BLUE_DARK, font_size=12, bold=True):
    add_rect(slide, left, top, width, height, fill_color=bg)
    add_text(slide, text, left, top, width, height,
             font_size=font_size, bold=bold, color=fg,
             align=PP_ALIGN.CENTER)


def icon_bullet(slide, items, left, top, width, font_size=15,
                color=SLATE_700, icon="●", spacing=0.42):
    """带 icon 的列表"""
    for i, (ico, txt, bold_flag, clr) in enumerate(items):
        y = top + i * spacing
        add_text(slide, ico, left, y, 0.35, 0.4,
                 font_size=font_size, bold=True, color=clr or BLUE_DARK)
        add_text(slide, txt, left + 0.35, y, width - 0.35, 0.4,
                 font_size=font_size, bold=bold_flag, color=clr or color)


def simple_table(slide, headers, rows, left, top, col_widths,
                 row_height=0.42, header_bg=BLUE_DARK, header_fg=WHITE,
                 alt_bg=BLUE_LIGHT, font_size=13):
    """简单表格"""
    total_w = sum(col_widths)
    # 表头
    x = left
    for i, (h, w) in enumerate(zip(headers, col_widths)):
        add_rect(slide, x, top, w, row_height, fill_color=header_bg)
        add_text(slide, h, x + 0.05, top + 0.04, w - 0.1, row_height - 0.08,
                 font_size=font_size, bold=True, color=header_fg,
                 align=PP_ALIGN.CENTER)
        x += w
    # 数据行
    for ri, row in enumerate(rows):
        bg = alt_bg if ri % 2 == 0 else WHITE
        x = left
        y = top + (ri + 1) * row_height
        for ci, (cell, w) in enumerate(zip(row, col_widths)):
            add_rect(slide, x, y, w, row_height, fill_color=bg,
                     line_color=SLATE_200, line_width=0.5)
            add_text(slide, str(cell), x + 0.06, y + 0.04,
                     w - 0.12, row_height - 0.08,
                     font_size=font_size, color=SLATE_900,
                     align=PP_ALIGN.CENTER)
            x += w


def box_with_title(slide, title, lines, left, top, width, height,
                   title_bg=BLUE_DARK, title_fg=WHITE,
                   body_bg=BLUE_LIGHT, font_size=14):
    title_h = 0.4
    add_rect(slide, left, top, width, title_h, fill_color=title_bg)
    add_text(slide, title, left + 0.1, top + 0.05, width - 0.2, title_h - 0.1,
             font_size=font_size, bold=True, color=title_fg)
    body_h = height - title_h
    add_rect(slide, left, top + title_h, width, body_h, fill_color=body_bg)
    body_items = [(ln,) if isinstance(ln, str) else ln for ln in lines]
    y = top + title_h + 0.1
    for item in body_items:
        txt = item[0]
        bold = item[1] if len(item) > 1 else False
        clr  = item[2] if len(item) > 2 else SLATE_900
        sz   = item[3] if len(item) > 3 else font_size - 1
        add_text(slide, txt, left + 0.15, y, width - 0.3, 0.35,
                 font_size=sz, bold=bold, color=clr)
        y += 0.36


# ══════════════════════════════════════════════════════════
# 各页幻灯片生成函数
# ══════════════════════════════════════════════════════════

def slide_01_cover(prs):
    """封面"""
    s = blank_slide(prs)
    # 全背景深蓝
    add_rect(s, 0, 0, 13.33, 7.5, fill_color=BLUE_DARK)
    # 右侧装饰条
    add_rect(s, 11.0, 0, 2.33, 7.5, fill_color=BLUE_MID)
    # 白色大标题
    add_text(s, "AutoTestDesign", 0.7, 1.6, 10, 1.2,
             font_size=52, bold=True, color=WHITE)
    add_text(s, "期末项目启动会", 0.7, 2.85, 10, 0.8,
             font_size=32, bold=False, color=RGBColor(0xBF, 0xDB, 0xFF))
    # 分隔线
    add_rect(s, 0.7, 3.75, 6, 0.06, fill_color=RGBColor(0x60, 0xA5, 0xFA))
    # 副标题
    add_text(s, "软件测试课 · Assignment 2 Final Project", 0.7, 3.95, 10, 0.5,
             font_size=18, color=RGBColor(0x93, 0xC5, 0xFD))
    add_text(s, "AI-Driven Test Design Tool for VCU Wake-Sleep Module",
             0.7, 4.55, 10, 0.5, font_size=15,
             color=RGBColor(0x7D, 0xD3, 0xFC))
    add_text(s, "2026 年 5 月", 0.7, 6.6, 4, 0.4,
             font_size=14, color=RGBColor(0xBF, 0xDB, 0xFF))


def slide_02_conclusion(prs):
    s = blank_slide(prs)
    header_bar(s, "先说结论", "我们的底牌是什么")
    footer_bar(s)

    add_rect(s, 0.5, 1.35, 12.33, 1.3, fill_color=BLUE_LIGHT)
    add_text(s, "💡  我们手里有一个已经跑通的完整系统：前端界面 + 后端 API + AI 模型",
             0.7, 1.45, 12, 0.5, font_size=17, bold=True, color=BLUE_DARK)
    add_text(s, "作业要我们做一个「AI驱动的自动测试工具」——我们要做的是改造和扩展，不是从零开始",
             0.7, 1.85, 12, 0.45, font_size=15, color=SLATE_700)

    items = [
        ("●", "现有代码量", True, BLUE_DARK),
        ("", "前端 ~40 个组件，后端 ~10 个 API 模块，全部可运行", False, SLATE_700),
        ("●", "真实数据", True, BLUE_DARK),
        ("", "5 个真实 BAIC VCU 测试数据库，共 9615 条测试记录", False, SLATE_700),
        ("●", "工作量评估", True, BLUE_DARK),
        ("", "5 人分工，约 4 周完成，每人约 30 小时", False, SLATE_700),
    ]
    icon_bullet(s, items, 0.7, 2.85, 8, font_size=15, spacing=0.45)

    # 右侧数字卡片
    for i, (num, label, sub) in enumerate([
        ("40+", "现有前端组件", "全部保留复用"),
        ("10+", "现有后端接口", "扩展新增模块"),
        ("9615", "真实测试记录", "5个数据库"),
    ]):
        x = 9.4
        y = 2.8 + i * 1.4
        add_rect(s, x, y, 3.4, 1.2, fill_color=BLUE_DARK)
        add_text(s, num, x, y + 0.1, 3.4, 0.55,
                 font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, label, x, y + 0.6, 3.4, 0.3,
                 font_size=13, bold=True, color=RGBColor(0xBF, 0xDB, 0xFF),
                 align=PP_ALIGN.CENTER)
        add_text(s, sub, x, y + 0.88, 3.4, 0.28,
                 font_size=11, color=RGBColor(0x93, 0xC5, 0xFD),
                 align=PP_ALIGN.CENTER)


def slide_03_what_is_system(prs):
    s = blank_slide(prs)
    header_bar(s, "现有项目是干什么的？", "用一句话和一个类比理解")
    footer_bar(s)

    # 左侧类比
    add_rect(s, 0.4, 1.3, 6.1, 5.6, fill_color=BLUE_LIGHT)
    add_text(s, "🔌  生活类比", 0.6, 1.45, 5.8, 0.45,
             font_size=17, bold=True, color=BLUE_DARK)
    lines = [
        "电动汽车插上充电枪",
        "↓",
        "汽车从睡眠状态被唤醒",
        "↓",
        "这个过程由 VCU（整车控制器）控制",
        "",
        "我们的系统：",
        "自动发送各种充电信号给 VCU",
        "观察 VCU 能不能正确响应",
    ]
    y = 2.05
    for line in lines:
        bold = line.startswith("我们") or line.startswith("自动")
        clr  = BLUE_DARK if bold else SLATE_700
        add_text(s, line, 0.65, y, 5.6, 0.38,
                 font_size=14, bold=bold, color=clr)
        y += 0.4

    # 右侧技术说明
    add_rect(s, 6.9, 1.3, 6.0, 5.6, fill_color=RGBColor(0xF8, 0xFA, 0xFC))
    add_text(s, "🛠  技术说明", 7.1, 1.45, 5.6, 0.45,
             font_size=17, bold=True, color=BLUE_DARK)

    tech_lines = [
        ("输入：向 VCU 发送测试信号", True, BLUE_DARK, 14),
        ("  CC2 电压值（如 6.3V / 4.7V）", False, SLATE_700, 13),
        ("  CC 电压值、CP 幅值、供电电压…", False, SLATE_700, 13),
        ("", False, SLATE_700, 10),
        ("输出：VCU 返回状态向量", True, BLUE_DARK, 14),
        ("  整车State状态：170（唤醒）/ 30（休眠）", False, SLATE_700, 13),
        ("  整车模式：5（唤醒）/ 2（休眠）", False, SLATE_700, 13),
        ("  READY 标志位：1 / 0", False, SLATE_700, 13),
        ("", False, SLATE_700, 10),
        ("结论：这就是软件黑盒测试", True, GREEN, 14),
        ("  给不同输入 → 检查输出是否符合预期", False, SLATE_700, 13),
    ]
    y = 2.05
    for item in tech_lines:
        txt, b, c, sz = item
        add_text(s, txt, 7.1, y, 5.6, 0.35, font_size=sz, bold=b, color=c)
        y += 0.38


def slide_04_vcu(prs):
    s = blank_slide(prs)
    header_bar(s, "VCU 是什么？", "3 分钟科普——不用记，看懂即可")
    footer_bar(s)

    # 左：状态图
    add_rect(s, 0.4, 1.3, 5.8, 5.6, fill_color=BLUE_LIGHT)
    add_text(s, "VCU 状态流程", 0.6, 1.38, 5.5, 0.4,
             font_size=16, bold=True, color=BLUE_DARK)

    flow = [
        ("充电桩", "输出 CC2 电压信号", BLUE_DARK, WHITE),
        ("VCU 整车控制器", "判断电压是否在有效范围", BLUE_MID, WHITE),
        ("唤醒 ✅  state=170", "CC2 在 [4.8, 7.7]V", GREEN, WHITE),
        ("休眠 💤  state=30", "CC2 = 12.0V", SLATE_700, WHITE),
        ("故障 ❌  state=30", "CC2 越界（<4.8 或 >7.8）", RED, WHITE),
    ]
    y = 1.95
    for title, sub, bg, fg in flow:
        add_rect(s, 0.7, y, 5.2, 0.55, fill_color=bg)
        add_text(s, title, 0.85, y + 0.03, 3.5, 0.28,
                 font_size=13, bold=True, color=fg)
        add_text(s, sub, 0.85, y + 0.27, 4.8, 0.25,
                 font_size=11, color=RGBColor(0xDB, 0xEA, 0xFE) if bg != GREEN else WHITE)
        y += 0.68
        if y < 6.4:
            add_text(s, "↓", 2.9, y - 0.18, 0.5, 0.25,
                     font_size=14, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)

    # 右：5个信号表格
    add_text(s, "5 种输入信号（我们要测的）", 6.6, 1.3, 6.5, 0.4,
             font_size=16, bold=True, color=BLUE_DARK)
    headers = ["信号名称", "物理含义", "有效区间"]
    rows = [
        ["CC2电压 ⭐", "充电唤醒主信号", "4.8~7.7V"],
        ["CC电压值", "充电枪接触检测", "> 4.0V"],
        ["CP幅值", "充电协议握手", "0V（待机）"],
        ["供电电压", "外部交流供电", "0V（无外供）"],
        ["网络唤醒", "远程唤醒指令", "0（关闭）"],
    ]
    simple_table(s, headers, rows, 6.5, 1.8,
                 [2.2, 2.3, 2.0], row_height=0.44, font_size=12)
    add_text(s, "⭐ CC2电压是最关键的主唤醒信号，其他4个是辅助检测信号",
             6.5, 5.5, 6.5, 0.4, font_size=12, color=ORANGE, bold=True)
    add_text(s, "所有边界值均来自 5 个真实 BAIC VCU 数据库（9615 条记录）分析",
             6.5, 5.95, 6.5, 0.4, font_size=11, color=SLATE_500)


def slide_05_gan(prs):
    s = blank_slide(prs)
    header_bar(s, "GAN 是什么？", "不用深究——知道它能干什么就够了")
    footer_bar(s)

    # 左：类比
    add_rect(s, 0.4, 1.3, 6.1, 5.6, fill_color=BLUE_LIGHT)
    add_text(s, "🎨  通俗类比：临摹大师", 0.6, 1.4, 5.7, 0.4,
             font_size=16, bold=True, color=BLUE_DARK)
    analogy = [
        "GAN = 造假大师 + 鉴别专家 对抗训练",
        "",
        "1. 造假大师 生成一张假名画",
        '2. 鉴别专家 说这里不像真的',
        "3. 造假大师 根据反馈改进",
        "4. 反复训练 → 以假乱真",
        "",
        "在我们项目里：",
        "• 真实数据 = 真实 CC2 电压序列",
        "• 造假大师 = GAN 生成器",
        "• 鉴别专家 = GAN 判别器",
        "• 结果 = 生成逼真的新电压序列",
    ]
    y = 1.95
    for line in analogy:
        bold = line.startswith("在我们") or "=" in line
        clr  = BLUE_DARK if bold else SLATE_700
        add_text(s, line, 0.65, y, 5.6, 0.35, font_size=13, bold=bold, color=clr)
        y += 0.36

    # 右：具体用途
    add_rect(s, 6.9, 1.3, 6.0, 5.6, fill_color=RGBColor(0xF8, 0xFA, 0xFC))
    add_text(s, "✅  GAN 在本项目中的角色", 7.1, 1.4, 5.6, 0.4,
             font_size=16, bold=True, color=BLUE_DARK)

    add_text(s, "EP / BVA 覆盖的场景（静态）：", 7.1, 1.95, 5.6, 0.35,
             font_size=13, bold=True, color=SLATE_900)
    add_text(s, "单次注入一个离散电压值\n→ 检验该值对应的输出状态", 7.3, 2.3, 5.4, 0.55,
             font_size=13, color=SLATE_700)

    add_text(s, "GAN 补充的场景（动态）：", 7.1, 3.05, 5.6, 0.35,
             font_size=13, bold=True, color=BLUE_DARK)
    add_text(s, "注入连续变化的 8 步电压序列\n→ 检验 VCU 在动态过程中\n   是否在正确位置发生状态转变",
             7.3, 3.4, 5.4, 0.75, font_size=13, color=SLATE_700)

    add_rect(s, 7.0, 4.35, 5.8, 0.55, fill_color=BLUE_LIGHT)
    add_text(s, "示例序列：[5.2, 6.1, 7.3, 7.6, 7.7, 7.8, 8.1]",
             7.1, 4.42, 5.6, 0.4, font_size=12, bold=True, color=BLUE_DARK)

    add_rect(s, 7.0, 5.1, 5.8, 0.9, fill_color=RGBColor(0xFF, 0xF7, 0xED))
    add_text(s, "⚠  GAN 只用于 CC2 电压",
             7.1, 5.15, 5.6, 0.35, font_size=13, bold=True, color=ORANGE)
    add_text(s, "其他4个信号用简单的 EP+BVA 即可\n（它们在数据库里就是线性扫描）",
             7.1, 5.48, 5.6, 0.5, font_size=12, color=SLATE_700)


def slide_06_existing(prs):
    s = blank_slide(prs)
    header_bar(s, "现有系统长什么样？", "已有功能 vs 需要新增的功能")
    footer_bar(s)

    # 左：已有
    add_rect(s, 0.4, 1.3, 5.9, 5.6, fill_color=RGBColor(0xF0, 0xFD, 0xF4))
    add_text(s, "✅  已有功能（不用重写）", 0.6, 1.4, 5.5, 0.4,
             font_size=16, bold=True, color=GREEN)
    have = [
        "前端 React + TypeScript + shadcn/ui",
        "后端 FastAPI + SQLite 完整框架",
        "GAN 模型推理接口（/api/gan）",
        "测试计划管理（CRUD）",
        "测试任务执行与监控",
        "实时 WebSocket 监控",
        "报告生成（JSON/Markdown）",
        "完整的 shadcn/ui 组件库",
    ]
    y = 1.95
    for item in have:
        add_text(s, "  ✓  " + item, 0.6, y, 5.6, 0.38,
                 font_size=13, color=RGBColor(0x16, 0x6F, 0x3A))
        y += 0.42

    # 右：需要新增
    add_rect(s, 6.7, 1.3, 6.2, 5.6, fill_color=RGBColor(0xFF, 0xF1, 0xF1))
    add_text(s, "🔨  需要新增的功能", 6.9, 1.4, 5.8, 0.4,
             font_size=16, bold=True, color=RED)
    need = [
        ("VCU 行为仿真器（目标应用）", True),
        ("需求导入与解析  FR 1.0/1.1", False),
        ("风险分析打分  FR 2.0", False),
        ("EP/BVA/决策表生成  FR 3.0", False),
        ("JSON/CSV/Excel 导出  FR 6.0", False),
        ("前端：4 个新功能页面", False),
        ("pytest 测试脚本与执行", False),
    ]
    y = 1.95
    for item, bold in need:
        clr = RED if bold else RGBColor(0x9B, 0x1C, 0x1C)
        add_text(s, "  ✗  " + item, 6.9, y, 5.8, 0.4,
                 font_size=13, bold=bold, color=clr)
        y += 0.5

    add_text(s, "合计新增工作量：约 70h 代码 + 80h 文档 + 25h 答辩",
             0.4, 6.75, 12.5, 0.35, font_size=12, color=SLATE_500,
             align=PP_ALIGN.CENTER)


def slide_07_requirements(prs):
    s = blank_slide(prs)
    header_bar(s, "Assignment 2 要求什么？", "关键理解：两个系统的关系")
    footer_bar(s)

    # 两个系统框图
    add_rect(s, 0.5, 1.35, 12.3, 2.0, fill_color=BLUE_LIGHT)
    add_text(s, "AutoTestDesign 工具（我们开发）", 0.8, 1.4, 6, 0.4,
             font_size=15, bold=True, color=BLUE_DARK)
    add_text(s, "负责：需求输入 → 解析 → 风险分析 → 用例生成 → 执行 → 导出",
             0.8, 1.78, 11, 0.38, font_size=13, color=SLATE_700)
    add_text(s, "↓  调用 API 测试", 5.5, 2.3, 2.5, 0.4,
             font_size=14, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)
    add_rect(s, 0.5, 2.65, 12.3, 0.65, fill_color=BLUE_DARK)
    add_text(s, "VCU 行为仿真器（目标应用，需要新建）",
             0.8, 2.73, 8, 0.4, font_size=15, bold=True, color=WHITE)
    add_text(s, "接收5种信号 → 返回 vehicle_state / test_status / ready_flag",
             0.8, 2.95, 11, 0.3, font_size=12, color=RGBColor(0xBF, 0xDB, 0xFF))

    # FR 表格
    add_text(s, "7 个功能需求", 0.5, 3.5, 4, 0.38,
             font_size=15, bold=True, color=BLUE_DARK)
    headers = ["编号", "功能", "是否必须"]
    rows = [
        ["FR 1.0", "需求导入（CSV/文本/表单）", "✅ 必须"],
        ["FR 1.1", "需求结构化解析", "✅ 必须"],
        ["FR 2.0", "风险分析与优先级", "✅ 必须"],
        ["FR 3.0", "黑盒用例生成（EP/BVA/决策表）", "✅ 必须"],
        ["FR 4.0", "白盒状态机建模", "⭐ 加分"],
        ["FR 6.0", "JSON/CSV/Excel 导出", "✅ 必须"],
        ["FR 7.0", "测试套件优化", "⭐ 加分"],
    ]
    simple_table(s, headers, rows, 0.5, 3.9,
                 [1.3, 6.3, 1.8], row_height=0.4, font_size=12)


def slide_08_scoring(prs):
    s = blank_slide(prs)
    header_bar(s, "评分权重", "交什么、怎么评——这是重中之重")
    footer_bar(s)

    # 左：4个交付物
    add_text(s, "四个交付物", 0.5, 1.35, 6, 0.4,
             font_size=16, bold=True, color=BLUE_DARK)
    deliverables = [
        ("20%", "工具本体", "代码 + README + 视频 Demo", BLUE_DARK),
        ("10%", "风险分析报告", "针对 VCU 仿真器的风险分析", SLATE_700),
        ("40%", "测试计划 Test Plan", "7个章节，最重要！", RED),
        ("30%", "详细测试设计与执行", "用例设计 + 脚本 + 结果分析", ORANGE),
    ]
    y = 1.85
    for pct, title, sub, clr in deliverables:
        add_rect(s, 0.5, y, 1.3, 0.85, fill_color=clr)
        add_text(s, pct, 0.5, y + 0.18, 1.3, 0.5,
                 font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, title, 1.95, y + 0.1, 4.5, 0.35,
                 font_size=14, bold=True, color=clr)
        add_text(s, sub, 1.95, y + 0.45, 4.5, 0.3,
                 font_size=12, color=SLATE_500)
        y += 1.1

    # 右：评分维度
    add_text(s, "五个评分维度", 7.2, 1.35, 6, 0.4,
             font_size=16, bold=True, color=BLUE_DARK)
    criteria = [
        ("10%", "概念理解", "ISO 29119-4，ISTQB，EP/BVA 是什么"),
        ("20%", "设计连贯性", "需求→风险→用例→执行，逻辑自洽"),
        ("40%", "覆盖度与有效性", "FR 全部实现，工具真的能用"),
        ("20%", "深度分析", "为什么这么设计，能推广吗"),
        ("10%", "答辩表现", "15分钟讲清楚，含 Q&A"),
    ]
    y = 1.85
    for pct, title, sub in criteria:
        add_rect(s, 7.2, y, 1.1, 0.75, fill_color=BLUE_DARK)
        add_text(s, pct, 7.2, y + 0.15, 1.1, 0.4,
                 font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, title, 8.45, y + 0.06, 4.6, 0.3,
                 font_size=14, bold=True, color=BLUE_DARK)
        add_text(s, sub, 8.45, y + 0.38, 4.6, 0.28,
                 font_size=12, color=SLATE_500)
        y += 0.95

    add_rect(s, 0.5, 6.55, 12.3, 0.45, fill_color=RGBColor(0xFF, 0xF7, 0xED))
    add_text(s, "⚠  Test Plan(40%) + Detailed Test Design(30%) = 70%   代码好不如文档好",
             0.7, 6.6, 12, 0.38, font_size=14, bold=True, color=ORANGE,
             align=PP_ALIGN.CENTER)


def slide_09_architecture(prs):
    s = blank_slide(prs)
    header_bar(s, "我们的方案：整体架构", "两个服务，一套完整流程")
    footer_bar(s)

    # 工具框
    add_rect(s, 0.4, 1.3, 12.5, 3.2, fill_color=BLUE_LIGHT)
    add_text(s, "AutoTestDesign Tool  （端口 8000）", 0.7, 1.38, 8, 0.42,
             font_size=17, bold=True, color=BLUE_DARK)

    boxes = [
        ("需求导入\nFR 1.0", "1.4"),
        ("需求解析\nFR 1.1", "3.1"),
        ("风险分析\nFR 2.0", "4.8"),
        ("用例生成\nFR 3.0+GAN", "6.5"),
        ("执行验证\nFR 5.0", "8.2"),
        ("导出报告\nFR 6.0", "9.9"),
    ]
    for label, x in boxes:
        add_rect(s, float(x), 1.95, 1.55, 1.1, fill_color=BLUE_DARK)
        add_text(s, label, float(x) + 0.05, 2.05, 1.45, 0.9,
                 font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # 箭头文字
    for x in [2.95, 4.65, 6.35, 8.05, 9.75]:
        add_text(s, "→", x, 2.38, 0.4, 0.35,
                 font_size=18, bold=True, color=BLUE_MID, align=PP_ALIGN.CENTER)

    add_rect(s, 0.4, 3.2, 12.5, 0.6, fill_color=BLUE_MID)
    add_text(s, "前端 React（端口 3000）  ↔  后端 FastAPI（端口 8000）  ↔  GAN 模型（CC2 电压时序）",
             0.7, 3.3, 12, 0.38, font_size=13, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(s, "↓  HTTP POST /simulate  发送单信号测试请求", 4.5, 3.95, 5, 0.45,
             font_size=14, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)

    # 仿真器框
    add_rect(s, 0.4, 4.55, 12.5, 2.0, fill_color=BLUE_DARK)
    add_text(s, "VCU 行为仿真器  （目标应用，端口 8001）",
             0.7, 4.65, 8, 0.42, font_size=17, bold=True, color=WHITE)
    add_text(s, "接收：5种输入信号之一（名称 + 值）",
             0.8, 5.15, 6, 0.35, font_size=13, color=RGBColor(0xBF, 0xDB, 0xFF))
    add_text(s, "返回：test_status  /  vehicle_state  /  vehicle_mode  /  ready_flag",
             0.8, 5.53, 11, 0.35, font_size=13, color=RGBColor(0xBF, 0xDB, 0xFF))
    add_text(s, "基于 9615 条真实记录建立的仿真逻辑",
             9.0, 5.15, 3.7, 0.35, font_size=12,
             color=RGBColor(0x7D, 0xD3, 0xFC), align=PP_ALIGN.RIGHT)


def slide_10_simulator(prs):
    s = blank_slide(prs)
    header_bar(s, "目标应用：VCU 行为仿真器", "基于真实数据库建立——老师了解项目背景")
    footer_bar(s)

    add_text(s, "为什么要自建仿真器？", 0.5, 1.35, 7, 0.38,
             font_size=15, bold=True, color=BLUE_DARK)
    add_text(s, "真实 BAIC HIL 台架目前无法连接 → 行业标准本来就是「先SIL仿真，再HIL测试」",
             0.5, 1.75, 12.3, 0.38, font_size=13, color=SLATE_700)

    # 仿真器响应表
    add_text(s, "仿真器核心响应规则（均来自真实数据库验证）：",
             0.5, 2.25, 8, 0.38, font_size=15, bold=True, color=BLUE_DARK)
    headers = ["发送信号", "测试值", "返回 vehicle_state", "含义"]
    rows = [
        ["CC2电压 ⭐", "6.3V（有效范围内）", "170", "唤醒成功 ✅"],
        ["CC2电压 ⭐", "9.0V（越界）", "30", "唤醒失败 ❌"],
        ["CC2电压 ⭐", "12.0V（休眠触发）", "30（status=3）", "正常休眠 💤"],
        ["CC电压值", "2.0V（接触不良区间）", "30", "检测失败 ❌"],
        ["CP幅值", "11.0V（协议冲突区间）", "30", "检测失败 ❌"],
        ["供电电压", "12.5V（过压区间）", "30", "检测失败 ❌"],
        ["网络唤醒", "1（冲突使能）", "30", "检测失败 ❌"],
    ]
    simple_table(s, headers, rows, 0.5, 2.7,
                 [2.5, 2.8, 2.6, 2.4], row_height=0.43, font_size=12)
    add_text(s, "3 个重要边界发现（来自数据分析，答辩会被问）：", 0.5, 6.35, 9, 0.35,
             font_size=12, bold=True, color=ORANGE)
    add_text(s, "① CC2=4.8V 能成功唤醒（db_2 6条确认）  ② 休眠测试5个信号全固定  ③ db_15 上界扩展到 8.1V",
             0.5, 6.7, 12.3, 0.35, font_size=11, color=SLATE_500)


def slide_11_testgen(prs):
    s = blank_slide(prs)
    header_bar(s, "测试用例怎么生成？", "三种黑盒技术 + GAN 补充")
    footer_bar(s)

    methods = [
        ("等价划分 EP", BLUE_DARK,
         "把输入范围分「有效类」和「无效类」\n每类取一个代表值",
         "CC2电压：\n  有效类取 6.3V\n  无效低取 3.0V\n  无效高取 9.0V"),
        ("边界值分析 BVA", RGBColor(0x7C, 0x3A, 0xED),
         "专测边界点\nbug 最容易藏在边界",
         "CC2 边界 4.8V：\n  4.7V（刚越界）\n  4.8V（精确边界）\n  4.9V（刚进入）"),
        ("决策表", GREEN,
         "测多信号组合条件\n一个信号异常→整体FAIL",
         "CC2有效 + CC电压值无效\n→ 整体 FAIL\n（来自真实测试策略）"),
        ("GAN 序列", ORANGE,
         "测电压连续变化过程\n静态单点测不到的场景",
         "[5.2, 6.1, 7.3, 7.6, 7.7,\n 7.8, 8.1, 9.0]\n→检测状态转变位置"),
    ]
    for i, (title, color, desc, example) in enumerate(methods):
        x = 0.4 + i * 3.22
        add_rect(s, x, 1.3, 3.0, 0.5, fill_color=color)
        add_text(s, title, x, 1.35, 3.0, 0.42,
                 font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(s, x, 1.8, 3.0, 2.0, fill_color=BLUE_LIGHT)
        add_text(s, desc, x + 0.1, 1.88, 2.8, 1.9,
                 font_size=12, color=SLATE_900)
        add_rect(s, x, 3.82, 3.0, 1.8, fill_color=RGBColor(0xF8, 0xFA, 0xFC))
        add_text(s, "示例：", x + 0.1, 3.88, 2.8, 0.3,
                 font_size=11, bold=True, color=color)
        add_text(s, example, x + 0.1, 4.15, 2.8, 1.4,
                 font_size=11, color=SLATE_700)

    # 汇总表
    add_text(s, "各信号使用的生成方法：", 0.4, 5.75, 6, 0.35,
             font_size=14, bold=True, color=BLUE_DARK)
    headers = ["信号", "EP", "BVA", "GAN"]
    rows = [
        ["CC2电压 ⭐", "✅", "✅", "✅时序序列"],
        ["CC电压值", "✅", "✅", "❌"],
        ["CP幅值", "✅", "✅", "❌"],
        ["供电电压", "✅", "✅", "❌"],
        ["网络唤醒", "✅穷举", "❌", "❌"],
    ]
    simple_table(s, headers, rows, 0.4, 6.15,
                 [2.5, 1.3, 1.3, 2.3], row_height=0.36, font_size=11)


def slide_12_workflow(prs):
    s = blank_slide(prs)
    header_bar(s, "工具的完整流程", "从需求到结果——满足作业的Mainly 要求")
    footer_bar(s)

    steps = [
        ("1", "导入 VCU 需求", "CSV / 文本 / 手动表单", "FR 1.0"),
        ("2", "结构化解析", "识别信号名/数值范围/条件/动作", "FR 1.1"),
        ("3", "风险分析打分", "5维度评分 → High/Medium/Low", "FR 2.0"),
        ("4", "生成测试用例", "EP + BVA + 决策表 + GAN 序列", "FR 3.0"),
        ("5", "调用仿真器执行", "发送信号 → 获取 vehicle_state", "FR 5.0"),
        ("6", "结果分析", "PASS/FAIL 统计，覆盖率报告", "分析"),
        ("7", "导出制品", "JSON / CSV / Excel 含追踪矩阵", "FR 6.0"),
    ]
    x_start = 0.35
    col_w = 12.5 / 7
    y_box = 1.35
    for i, (num, title, sub, tag_txt) in enumerate(steps):
        x = x_start + i * col_w
        add_rect(s, x, y_box, col_w - 0.1, 0.6, fill_color=BLUE_DARK)
        add_text(s, num, x, y_box + 0.1, col_w - 0.1, 0.42,
                 font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(s, x, y_box + 0.6, col_w - 0.1, 1.1, fill_color=BLUE_LIGHT)
        add_text(s, title, x + 0.05, y_box + 0.68, col_w - 0.2, 0.45,
                 font_size=12, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)
        add_text(s, sub, x + 0.05, y_box + 1.1, col_w - 0.2, 0.55,
                 font_size=10, color=SLATE_700, align=PP_ALIGN.CENTER)
        tag(s, tag_txt, x + 0.05, y_box + 1.7, col_w - 0.2, 0.3,
            font_size=9)
        if i < 6:
            add_text(s, "→", x + col_w - 0.2, y_box + 0.9, 0.3, 0.4,
                     font_size=14, bold=True, color=BLUE_MID, align=PP_ALIGN.CENTER)

    add_rect(s, 0.4, 2.25, 12.5, 0.65,
             fill_color=RGBColor(0xF0, 0xFD, 0xF4))
    add_text(s, "🔄  Interactive Review（强制要求）：每个步骤用户都可以介入修改，不是可选功能！",
             0.7, 2.35, 12, 0.4, font_size=13, bold=True, color=GREEN)

    # 下半部分说明
    add_text(s, "每个阶段的可交互内容：", 0.4, 3.15, 6, 0.38,
             font_size=14, bold=True, color=BLUE_DARK)
    review_items = [
        "需求导入 → 内联编辑任何需求字段",
        "解析结果 → 增删改 Input Fields / Data Ranges / Conditions",
        "风险分析 → 手动覆盖任意维度分值和优先级",
        "测试用例 → 编辑/删除/手动新增任意用例",
        "导出内容 → 勾选要导出的模块和格式",
    ]
    y = 3.6
    for item in review_items:
        add_text(s, "  ✦  " + item, 0.5, y, 12, 0.38,
                 font_size=13, color=SLATE_700)
        y += 0.42


def slide_13_data(prs):
    s = blank_slide(prs)
    header_bar(s, "所有设计都有真实数据支撑", "5 个数据库，9615 条记录分析")
    footer_bar(s)

    # 数据库概览
    add_text(s, "我们分析的数据：", 0.5, 1.35, 6, 0.38,
             font_size=15, bold=True, color=BLUE_DARK)
    dbs = [("db_10.db", "1010 条"), ("db_11.db", "1000 条"),
           ("db_15.db", "1488 条"), ("db.db", "3000 条"), ("db_2.db", "1617 条")]
    x = 0.5
    for name, cnt in dbs:
        add_rect(s, x, 1.78, 2.3, 0.75, fill_color=BLUE_DARK)
        add_text(s, name, x + 0.05, 1.84, 2.2, 0.35,
                 font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, cnt, x + 0.05, 2.16, 2.2, 0.3,
                 font_size=11, color=RGBColor(0xBF, 0xDB, 0xFF),
                 align=PP_ALIGN.CENTER)
        x += 2.45

    # 3个重要发现
    add_text(s, "3 个来自数据的重要发现（答辩必考）：",
             0.5, 2.75, 10, 0.38, font_size=15, bold=True, color=BLUE_DARK)
    findings = [
        ("发现 1", "CC2=4.8V 能成功唤醒 VCU",
         "db_2 中 6 条记录：4.8V → vehicle_state=170（成功唤醒）\n"
         "但数据库 test_status=4，因为测试框架期望它失败——实际 VCU 成功响应了",
         BLUE_DARK),
        ("发现 2", "休眠测试 5 个信号全部固定",
         "1500 条 type=2 记录：5 个信号值完全不变\n"
         "CC2=12.0V, CC=12.0V, CP=0, 供电=0, 网络唤醒=0\n"
         "变化来自 VCU 的当前状态，不来自输入信号",
         GREEN),
        ("发现 3", "db_15 CC2 有效上界扩展到 8.1V",
         "其余 4 个数据库：CC2 最高 PASS = 7.7V\n"
         "db_15：CC2 最高 PASS = 8.1V（939 条记录）\n"
         "代表不同 VCU 固件批次的配置差异",
         ORANGE),
    ]
    y = 3.28
    for i, (tag_txt, title, detail, clr) in enumerate(findings):
        x = 0.4 + i * 4.3
        add_rect(s, x, y, 4.0, 0.4, fill_color=clr)
        add_text(s, f"{tag_txt}：{title}", x + 0.1, y + 0.06, 3.8, 0.3,
                 font_size=13, bold=True, color=WHITE)
        add_rect(s, x, y + 0.4, 4.0, 2.6, fill_color=BLUE_LIGHT if i == 0 else (
            RGBColor(0xF0, 0xFD, 0xF4) if i == 1 else YELLOW_BG))
        add_text(s, detail, x + 0.1, y + 0.5, 3.8, 2.4,
                 font_size=11, color=SLATE_700)

    add_text(s, "这说明：真实系统行为比规格说明更复杂  ——  AutoTestDesign 工具的价值正在于系统化发现这些差异",
             0.4, 6.65, 12.5, 0.38, font_size=12, bold=True, color=BLUE_DARK,
             align=PP_ALIGN.CENTER)


def slide_14_timeline(prs):
    s = blank_slide(prs)
    header_bar(s, "4 周工作计划", "关键里程碑与每人每周目标")
    footer_bar(s)

    weeks = ["第 1 周", "第 2 周", "第 3 周", "第 4 周"]
    members = ["A 后端架构", "B 测试算法", "C 前端", "D 文档", "E 统筹"]
    tasks = [
        ["⭐ VCU 仿真器", "FR1.0/1.1 后端", "Test Plan技术节", "修改+支持"],
        ["对齐 schemas", "FR2.0/3.0 算法", "pytest 执行+文档", "修改+支持"],
        ["Mock 数据开发", "4页面开发", "联调+截图", "README+视频"],
        ["Risk Report草稿", "Test Plan 主体", "文档收尾", "配合修改"],
        ["PPT模板+Q&A", "Test Plan 汇总", "⭐ PPT+校对", "⭐ 答辩排练"],
    ]
    colors = [BLUE_DARK, RGBColor(0x7C, 0x3A, 0xED), GREEN, ORANGE, SLATE_700]

    col_w = 2.55
    row_h = 0.85
    header_h = 0.4

    # 周次表头
    for j, week in enumerate(weeks):
        x = 1.95 + j * col_w
        add_rect(s, x, 1.3, col_w - 0.1, header_h, fill_color=BLUE_DARK)
        add_text(s, week, x, 1.35, col_w - 0.1, header_h - 0.1,
                 font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # 成员行
    for i, (member, task_row, clr) in enumerate(zip(members, tasks, colors)):
        y = 1.3 + header_h + i * row_h
        add_rect(s, 0.4, y, 1.5, row_h - 0.08, fill_color=clr)
        add_text(s, member, 0.4, y + 0.1, 1.5, row_h - 0.28,
                 font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        for j, task in enumerate(task_row):
            x = 1.95 + j * col_w
            bg = BLUE_LIGHT if (i + j) % 2 == 0 else WHITE
            add_rect(s, x, y, col_w - 0.1, row_h - 0.08,
                     fill_color=bg, line_color=SLATE_200, line_width=0.5)
            add_text(s, task, x + 0.08, y + 0.12, col_w - 0.25, row_h - 0.28,
                     font_size=11, color=SLATE_900 if "⭐" not in task else clr,
                     bold="⭐" in task)

    # 里程碑
    milestones = [
        ("第1周末", "仿真器跑通 → 其他人开始对接"),
        ("第2周末", "前端+后端联调"),
        ("第3周末", "代码完成 + 文档草稿完成"),
        ("第4周初", "PPT + 视频 → 提交 + 答辩"),
    ]
    y = 6.5
    for j, (week, text) in enumerate(milestones):
        x = 0.4 + j * 3.2
        add_rect(s, x, y, 0.9, 0.5, fill_color=BLUE_DARK)
        add_text(s, week, x, y + 0.08, 0.9, 0.35,
                 font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, text, x + 1.0, y + 0.1, 2.1, 0.35,
                 font_size=11, color=SLATE_700)


def slide_15_roles_overview(prs):
    s = blank_slide(prs)
    header_bar(s, "五个人，五个角色", "根据自己的技能和意愿选择")
    footer_bar(s)

    roles = [
        ("A", "后端架构师", "Python/FastAPI", "搭系统骨架，最早完成，其他人都等你", BLUE_DARK),
        ("B", "测试算法工程师", "Python + 逻辑", "写工具最核心的算法，同时产出 Detailed Test Design", RGBColor(0x7C, 0x3A, 0xED)),
        ("C", "前端工程师", "React / TypeScript", "做用户看得见的界面，录制演示视频", GREEN),
        ("D", "测试分析师", "文档写作", "写风险分析报告和 Test Plan 主体，占总分 50%！", ORANGE),
        ("E", "统筹负责人", "统筹 + PPT", "汇总所有人工作，制作 PPT，组织答辩", SLATE_700),
    ]
    for i, (letter, name, skill, desc, clr) in enumerate(roles):
        y = 1.35 + i * 1.12
        add_rect(s, 0.4, y, 0.9, 0.9, fill_color=clr)
        add_text(s, letter, 0.4, y + 0.15, 0.9, 0.6,
                 font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(s, 1.4, y, 11.5, 0.9, fill_color=BLUE_LIGHT if i % 2 == 0 else WHITE,
                 line_color=SLATE_200, line_width=0.5)
        add_text(s, name, 1.6, y + 0.08, 3.5, 0.38,
                 font_size=16, bold=True, color=clr)
        tag(s, skill, 5.3, y + 0.1, 2.5, 0.3, bg=clr, fg=WHITE, font_size=11)
        add_text(s, desc, 1.6, y + 0.5, 11, 0.35,
                 font_size=12, color=SLATE_700)


def member_slide(prs, letter, name, color, code_tasks, doc_tasks, deliverables, note=""):
    s = blank_slide(prs)
    add_rect(s, 0, 0, 13.33, 1.15, fill_color=color)
    add_rect(s, 0, 0, 1.4, 1.15, fill_color=BLUE_MID)
    add_text(s, letter, 0, 0.15, 1.4, 0.85,
             font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, f"成员 {letter} — {name}", 1.5, 0.15, 9, 0.5,
             font_size=26, bold=True, color=WHITE)
    add_rect(s, 0, 1.15, 13.33, 0.04, fill_color=RGBColor(0x60, 0xA5, 0xFA))
    footer_bar(s)

    # 代码任务
    add_rect(s, 0.4, 1.3, 5.9, 0.4, fill_color=color)
    add_text(s, "代码任务", 0.5, 1.36, 5.7, 0.3,
             font_size=14, bold=True, color=WHITE)
    y = 1.8
    for task in code_tasks:
        add_text(s, "  ▸  " + task, 0.4, y, 5.9, 0.38,
                 font_size=13, color=SLATE_900)
        y += 0.42
    # 补满
    code_end = y

    # 文档任务
    add_rect(s, 6.7, 1.3, 5.9, 0.4, fill_color=SLATE_700)
    add_text(s, "文档 / 答辩任务", 6.8, 1.36, 5.7, 0.3,
             font_size=14, bold=True, color=WHITE)
    y = 1.8
    for task in doc_tasks:
        add_text(s, "  ▸  " + task, 6.7, y, 5.9, 0.38,
                 font_size=13, color=SLATE_900)
        y += 0.42

    # 交付物
    add_rect(s, 0.4, 5.7, 12.2, 0.35, fill_color=SLATE_200)
    add_text(s, "📦  关键交付物", 0.55, 5.73, 3, 0.28,
             font_size=12, bold=True, color=SLATE_700)
    deliv_text = "  |  ".join(deliverables)
    add_text(s, deliv_text, 3.4, 5.73, 9, 0.28,
             font_size=11, color=SLATE_700)

    if note:
        add_rect(s, 0.4, 6.15, 12.2, 0.55, fill_color=YELLOW_BG)
        add_text(s, "💡  " + note, 0.6, 6.22, 12, 0.42,
                 font_size=12, color=ORANGE, bold=True)


def slide_16_member_a(prs):
    member_slide(prs, "A", "后端架构师", BLUE_DARK,
        code_tasks=[
            "VCU 仿真器（vcu_simulator/，端口8001）⭐最高优先级",
            "扩展 schemas.py（与B对齐数据模型）",
            "FR 1.0/1.1：需求导入API + 正则解析服务",
            "main.py 注册所有新路由",
        ],
        doc_tasks=[
            "simulator_spec_for_doc.md → 给D写文档用",
            "API_reference_for_frontend.md → 给C开发用",
            "Test Plan: Testing Framework章节",
            "Test Plan: Cost Estimation章节",
            "Test Plan: Schedule / Checklist章节",
        ],
        deliverables=["第1周末：仿真器跑通", "两份说明文档", "第2周末：FR1.0/1.1完成"],
        note="你是关键路径！仿真器必须第1周末完成，否则B和C都无法开始联调"
    )


def slide_17_member_b(prs):
    member_slide(prs, "B", "测试算法工程师", RGBColor(0x7C, 0x3A, 0xED),
        code_tasks=[
            "FR 2.0：5维度风险评分算法",
            "FR 3.0：EP / BVA / 决策表生成器",
            "simulator_client.py（调用仿真器的HTTP客户端）",
            "FR 6.0：JSON / CSV / Excel 导出服务",
            "pytest 测试脚本（37条参数化用例）+ 实际执行",
        ],
        doc_tasks=[
            "detailed_test_case_design.md",
            "  → EP/BVA/决策表完整用例表格 + 选型理由",
            "test_result_analysis.md",
            "  → pytest结果统计 + 追踪矩阵 + 改进说明",
        ],
        deliverables=["第2周末：FR2/3/6后端完成", "第3周：pytest_results.txt", "两份文档"],
        note="你写代码的同时同步写文档——这两份文档直接变成 Detailed Test Design（30分）"
    )


def slide_18_member_c(prs):
    member_slide(prs, "C", "前端工程师", GREEN,
        code_tasks=[
            "修改 Sidebar.tsx / App.tsx / api.ts（第1周先做）",
            "RequirementInput.tsx（需求导入+解析展示）",
            "RiskAnalysis.tsx（风险矩阵+分值调整）",
            "TestCaseDesign.tsx（用例表格+执行进度）",
            "ExportCenter.tsx（格式选择+下载）",
        ],
        doc_tasks=[
            "每个页面实现 Interactive Review（内联编辑）",
            "录制视频 Demo（6分钟，按脚本）",
            "更新 README.md（新页面说明）",
            "整理各页面截图（约8~10张，给E做PPT）",
        ],
        deliverables=["第2周初：Sidebar/App修改完", "第3周末：4页面联调完+截图", "第4周：视频Demo"],
        note="等不到接口时用 Mock 数据先开发 UI，联调时再换真实接口"
    )


def slide_19_member_d(prs):
    member_slide(prs, "D", "测试分析师", ORANGE,
        code_tasks=[
            "（不是主力写代码）",
            "协助：运行 pytest 脚本获取测试截图",
            "整理风险矩阵图（工具截图 or Excel制作）",
        ],
        doc_tasks=[
            "Artifact 2（完整）：Risk Analysis Report",
            "  → 10条需求 × 5维度评分表 + 风险矩阵图",
            "Test Plan: Project Scope章节",
            "Test Plan: Test Items章节",
            "Test Plan: High-level Suite Design（7套件）",
        ],
        deliverables=["第2周末：Risk Report草稿", "第3周初：Test Plan主体三章节", "风险矩阵图"],
        note="你写的 Risk Analysis Report(10%) + Test Plan主体(40%) 合计决定50%的分值"
    )


def slide_20_member_e(prs):
    member_slide(prs, "E", "统筹负责人", SLATE_700,
        code_tasks=[
            "（不负责独立代码模块）",
            "负责整体联调协调",
            "检查各服务能否互相正确调用",
        ],
        doc_tasks=[
            "Test Plan: Organization Chart章节",
            "汇总 Test Plan（校对一致性）",
            "汇总 Detailed Test Design（补写Tool Implementation节）",
            "制作 PPT（18~22页）",
            "Q&A问题库 + 答辩脚本（15分钟分配）",
            "组织全员答辩排练（至少1次）",
        ],
        deliverables=["第3周末：Q&A问题库", "第4周初：PPT完成", "答辩脚本分发"],
        note="你需要理解所有人的工作——建议每周和每个成员同步一次进度"
    )


def slide_21_vote(prs):
    s = blank_slide(prs)
    header_bar(s, "现在来选！", "每人选一个角色，现场确认")
    footer_bar(s)

    headers = ["角色", "主要技能需求", "工作重心", "意向人选"]
    rows = [
        ["A — 后端架构师", "Python / FastAPI", "代码（仿真器+FR1）", ""],
        ["B — 测试算法工程师", "Python / 逻辑思维", "代码（FR2/3/6）+ 文档", ""],
        ["C — 前端工程师", "React / TypeScript", "代码（前端页面）+ 视频", ""],
        ["D — 测试分析师", "文档写作能力", "文档（Risk Report + Test Plan）", ""],
        ["E — 统筹负责人", "统筹能力 / PPT", "汇总整合 + 答辩", ""],
    ]
    simple_table(s, headers, rows, 0.4, 1.35,
                 [3.0, 2.8, 3.5, 2.8], row_height=0.65, font_size=13)

    add_text(s, "选择参考：", 0.4, 5.45, 4, 0.38,
             font_size=14, bold=True, color=BLUE_DARK)
    hints = [
        ("Python / 后端经验", "→ 优先选 A 或 B"),
        ("React / 前端经验", "→ 选 C"),
        ("文档能力强，不一定会写代码", "→ 选 D"),
        ("愿意统筹协调，喜欢答辩", "→ 选 E"),
    ]
    y = 5.9
    for skill, choice in hints:
        add_text(s, f"  •  {skill}", 0.4, y, 6, 0.35,
                 font_size=13, color=SLATE_700)
        add_text(s, choice, 6.5, y, 3, 0.35,
                 font_size=13, bold=True, color=BLUE_DARK)
        y += 0.42


def slide_22_action(prs):
    s = blank_slide(prs)
    header_bar(s, "会后行动清单", "今晚确定，明天开始")
    footer_bar(s)

    # 左：立刻要做
    add_rect(s, 0.4, 1.3, 6.0, 5.5, fill_color=BLUE_LIGHT)
    add_text(s, "📋  今晚确认", 0.6, 1.38, 5.6, 0.38,
             font_size=16, bold=True, color=BLUE_DARK)
    actions = [
        "□  5人分工表填写完毕",
        "□  每人阅读自己的任务书",
        "   docs/tasks/member_X_xxx.md",
        "□  A + B 约 schemas 对齐会（30分钟）",
        "□  E 建立共享文档空间",
        "   （飞书 / Notion / Google Docs）",
    ]
    y = 1.85
    for item in actions:
        bold = item.startswith("□")
        add_text(s, item, 0.6, y, 5.6, 0.38,
                 font_size=13, bold=bold, color=SLATE_900)
        y += 0.42

    # 右：本周重点
    add_rect(s, 6.8, 1.3, 6.1, 5.5, fill_color=RGBColor(0xF8, 0xFA, 0xFC))
    add_text(s, "🚀  本周重点（第1周）", 7.0, 1.38, 5.7, 0.38,
             font_size=16, bold=True, color=BLUE_DARK)
    week1 = [
        ("A", "开始写 VCU 仿真器", BLUE_DARK),
        ("B", "阅读 DESIGN_PLAN.md，与A对齐", RGBColor(0x7C, 0x3A, 0xED)),
        ("C", "搭前端框架，用Mock数据先开发", GREEN),
        ("D", "阅读 DESIGN_PLAN.md 第二节", ORANGE),
        ("E", "准备PPT模板，整理Q&A框架", SLATE_700),
    ]
    y = 1.85
    for letter, task, clr in week1:
        add_rect(s, 7.0, y, 0.6, 0.38, fill_color=clr)
        add_text(s, letter, 7.0, y + 0.04, 0.6, 0.3,
                 font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, task, 7.7, y, 5.0, 0.38,
                 font_size=13, color=SLATE_900)
        y += 0.5

    # 底部文件位置
    add_rect(s, 0.4, 6.55, 12.5, 0.55, fill_color=BLUE_DARK)
    add_text(s, "📁  docs/DESIGN_PLAN.md（必读）  |  docs/tasks/member_X.md（各自任务书）",
             0.6, 6.62, 12.2, 0.4, font_size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════
# 主函数
# ══════════════════════════════════════════════════════════

def main():
    prs = new_prs()

    print("生成幻灯片...")
    slide_01_cover(prs)        ; print("  P1  封面")
    slide_02_conclusion(prs)   ; print("  P2  先说结论")
    slide_03_what_is_system(prs);print("  P3  系统是干什么的")
    slide_04_vcu(prs)          ; print("  P4  VCU 科普")
    slide_05_gan(prs)          ; print("  P5  GAN 科普")
    slide_06_existing(prs)     ; print("  P6  现有系统")
    slide_07_requirements(prs) ; print("  P7  作业要求")
    slide_08_scoring(prs)      ; print("  P8  评分权重")
    slide_09_architecture(prs) ; print("  P9  整体架构")
    slide_10_simulator(prs)    ; print("  P10 目标应用")
    slide_11_testgen(prs)      ; print("  P11 测试生成")
    slide_12_workflow(prs)     ; print("  P12 工具流程")
    slide_13_data(prs)         ; print("  P13 数据支撑")
    slide_14_timeline(prs)     ; print("  P14 时间规划")
    slide_15_roles_overview(prs);print("  P15 角色概览")
    slide_16_member_a(prs)     ; print("  P16 成员 A")
    slide_17_member_b(prs)     ; print("  P17 成员 B")
    slide_18_member_c(prs)     ; print("  P18 成员 C")
    slide_19_member_d(prs)     ; print("  P19 成员 D")
    slide_20_member_e(prs)     ; print("  P20 成员 E")
    slide_21_vote(prs)         ; print("  P21 选择")
    slide_22_action(prs)       ; print("  P22 行动清单")

    out_path = "docs/AutoTestDesign_启动会.pptx"
    prs.save(out_path)
    print(f"\n✅  PPT 已生成：{out_path}")
    print(f"   共 {len(prs.slides)} 页")


if __name__ == "__main__":
    main()
